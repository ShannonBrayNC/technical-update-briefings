# Write a complete run_build.py implementing the features discussed.
from pathlib import Path
code = r'''#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
run_build.py
------------
Builds a styled PowerPoint deck from HTML inputs (Roadmap + Message Center).
- MC-first merge with fuzzy dedupe
- Right rail with product-colored background, big status chip, status icon
- Rail card (Month, Roadmap ID, compact details)
- "Target Audience" header + two rows (End Users, Admins)
- Left side: centered title with auto height clamp; summary paragraph (no bullets)
- Background comes from brand_bg (else cover)
- Footer shows month + page/total

CLI (underscored flags) for direct invocation.
Use generate_deck.py (hyphenated flags) if you need style.yaml support.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Tuple, Union
from dataclasses import dataclass

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ---- slides helpers (local) -------------------------------------------------
# Expect a local slides.py providing add_full_slide_picture/add_title_box/add_text_box.
# We import safely and fail clearly if missing.
import importlib.util as _impspec
_slides_mod = None
for cand in ("slides.py", "deck_slides.py"):
    p = Path(__file__).with_name(cand)
    if p.exists():
        spec = _impspec.spec_from_file_location("slides", str(p))
        if spec and spec.loader:
            _slides_mod = _impspec.module_from_spec(spec)
            spec.loader.exec_module(_slides_mod)
            break
if _slides_mod is None:
    raise ImportError("Local slides.py/deck_slides.py not found or failed to import.")
S = _slides_mod

# ---- logging ----------------------------------------------------------------
def _log(msg: str) -> None:
    print(f"[run_build] {msg}", flush=True)

# ---- simple HTML parsers (prefer external 'parsers' if available) -----------
def _try_external_parsers() -> Optional[object]:
    try:
        spec = _impspec.find_spec("parsers")
        if spec is None:
            return None
        mod = importlib.import_module("parsers")  # type: ignore
        return mod
    except Exception:
        return None

def _parse_one(path: str, month: str) -> List[Dict]:
    """Parse a single HTML file into a list of item dicts.
    Tries external 'parsers' first; else uses a forgiving built-in fallback.
    """
    mod = _try_external_parsers()
    p = Path(path)
    src = "message_center" if "message_center" in str(p).lower() else "roadmap"
    items: List[Dict] = []
    if mod:
        try:
            if src == "message_center" and hasattr(mod, "parse_message_center_html"):
                items = mod.parse_message_center_html(str(p), month)  # type: ignore
            elif src == "roadmap" and hasattr(mod, "parse_roadmap_html"):
                items = mod.parse_roadmap_html(str(p), month)  # type: ignore
        except Exception as e:
            _log(f"ERROR external parser on {path}: {e}")

    if items:
        _log(f"Parsed {len(items)} with month='{month}' from {path}")
        for it in items:
            it.setdefault("_source", src)
        return items

    # --- fallback (very naive) ---
    try:
        text = Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        _log(f"ERROR parsing {path}: {e}")
        return []

    # extract sections that look like items (very permissive)
    blocks = re.split(r"\n{2,}|<hr[^>]*>", text, flags=re.I)
    for b in blocks:
        title = ""
        m = re.search(r"(?:^|\n)\s*Title\s*:\s*(.+)", b, flags=re.I)
        if m:
            title = m.group(1).strip()
        else:
            # fallback: first non-empty line
            lines = [ln.strip() for ln in re.split(r"<[^>]+>|\r|\n", b) if ln.strip()]
            title = lines[0][:140] if lines else ""

        if not title:
            continue

        item: Dict = {
            "title": title,
            "summary": "",
            "description": "",
            "status": "",
            "roadmap_id": "",
            "products": [],
            "platforms": [],
            "clouds": [],
            "ga": "",
            "_source": src,
        }

        def grab(label: str) -> str:
            m = re.search(rf"{label}\s*:\s*(.+)", b, flags=re.I)
            return (m.group(1).strip() if m else "")

        item["status"] = grab("Status")
        item["roadmap_id"] = grab("Roadmap ID") or grab("ID")
        item["summary"] = grab("Summary") or grab("Description")
        item["description"] = item["summary"]
        prod = grab("Product") or grab("Products")
        if prod:
            item["products"] = [p.strip() for p in re.split(r",|/|;", prod) if p.strip()]
        plat = grab("Platform") or grab("Platforms")
        if plat:
            item["platforms"] = [p.strip() for p in re.split(r",|/|;", plat) if p.strip()]
        clouds = grab("Cloud") or grab("Clouds")
        if clouds:
            item["clouds"] = [p.strip() for p in re.split(r",|/|;", clouds) if p.strip()]
        item["ga"] = grab("GA") or grab("Release") or ""

        items.append(item)

    _log(f"Parsed {len(items)} with month='{month}' from {path}")
    return items

def _parse_inputs(inputs: List[str], month: str) -> List[Dict]:
    all_items: List[Dict] = []
    for p in inputs:
        all_items.extend(_parse_one(p, month))
    return all_items

# ---- merge & fuzzy dedupe (MC-first) ----------------------------------------
def _norm(s: str) -> str:
    s = (s or "").lower()
    s = re.sub(r"[^a-z0-9]+", " ", s)
    return re.sub(r"\s+", " ", s).strip()

def _token_set(s: str) -> set:
    return set(_norm(s).split()) if s else set()

def _similar(a: str, b: str) -> float:
    """Jaccard similarity over tokens of titles."""
    ta, tb = _token_set(a), _token_set(b)
    if not ta or not tb:
        return 0.0
    inter = len(ta & tb)
    union = len(ta | tb)
    return inter / union if union else 0.0

def _mc_first_merge(items: List[Dict]) -> List[Dict]:
    """Prefer Message Center items; fuzzy dedupe by title and roadmap_id."""
    out: List[Dict] = []
    seen_by_id: Dict[str, Dict] = {}
    THRESH = 0.62

    # First, split MC and RM
    mc = [it for it in items if it.get("_source") == "message_center"]
    rm = [it for it in items if it.get("_source") != "message_center"]

    # Index MC by roadmap_id/title
    for it in mc:
        rid = str(it.get("roadmap_id") or "").strip()
        title = it.get("title") or ""
        if rid:
            seen_by_id[rid] = it
        out.append(it)

    # Merge RM items
    for it in rm:
        rid = str(it.get("roadmap_id") or "").strip()
        title = it.get("title") or ""
        if rid and rid in seen_by_id:
            # already have MC item, skip
            continue
        # fuzzy against all MC titles
        sim = max((_similar(title, m.get("title") or "") for m in mc), default=0.0)
        if sim >= THRESH:
            continue
        out.append(it)

    return out

# ---- right rail rendering helpers -------------------------------------------
def _add_rail(slide, rail_width_in: float = 3.5, hex_color: str = "0F172A"):
    """Draw the right vertical rail across full height."""
    if not rail_width_in or rail_width_in <= 0:
        return
    hex_color = (hex_color or "0F172A").replace("#", "")
    full_w, full_h = int(10 * 914400), int(7.5 * 914400)
    rail_w = int(rail_width_in * 914400)
    left   = full_w - rail_w
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, 0, rail_w, full_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    shp.line.fill.background()

def _rail_color_for(subject: Union[Dict, Sequence[str], str, None], assets: Dict) -> str:
    """
    Pick a rail color from assets['product_palette'].
    Accepts either the whole item dict, a products list/tuple/set, a single product str, or None.
    """
    palette = assets.get("product_palette") or {}
    if not isinstance(palette, dict):
        palette = {}
    # normalize keys to lowercase
    if palette:
        for k in list(palette.keys()):
            if isinstance(k, str):
                lk = k.lower()
                if lk != k:
                    palette[lk] = palette.pop(k)

    def normalize_products(x) -> List[str]:
        if x is None:
            return []
        if isinstance(x, dict):
            x = x.get("products") or x.get("product") or []
        if isinstance(x, str):
            return [x]
        if isinstance(x, (list, tuple, set)):
            return [str(v) for v in x if v is not None]
        return []

    prods = normalize_products(subject)
    key = (prods[0] if prods else "").strip().lower()
    color = palette.get(key) or palette.get("default") or "0F172A"
    return str(color).replace("#", "")

def _status_icon_for(status: str, assets: Dict) -> Optional[str]:
    s = (status or "").strip().lower()
    rocket  = assets.get("icon_rocket")
    preview = assets.get("icon_preview")
    if any(k in s for k in ("launched", "rolling out", "ga", "general availability", "available", "rolled out")):
        return rocket if rocket and Path(rocket).exists() else None
    if any(k in s for k in ("in development", "development", "dev", "preview")):
        return preview if preview and Path(preview).exists() else None
    return None

def _status_icon(slide, rail_left_in: float, rail_w_in: float, item: Dict, assets: Dict) -> None:
    path = _status_icon_for(item.get("status"), assets)
    if not path:
        return
    ICON_W, ICON_H = 0.90, 0.90
    left = rail_left_in + rail_w_in - ICON_W - 0.35
    top  = 0.10
    try:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 width=Inches(ICON_W), height=Inches(ICON_H))
    except Exception:
        pass

def _status_chip(slide, rail_left_in: float, rail_w_in: float, status_text: str):
    """Draw a larger rounded 'In Development' / 'Rolling Out' chip centered in the rail."""
    if rail_w_in <= 0:
        return
    CHIP_W_IN, CHIP_H_IN = 3.0, 0.78
    CHIP_FONT_PT = 18
    left = rail_left_in + (rail_w_in - CHIP_W_IN) / 2.0
    top  = 0.60
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(CHIP_W_IN), Inches(CHIP_H_IN))
    shp.adjustments[0] = 0.2
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string("495057")
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = status_text
    r.font.size = Pt(CHIP_FONT_PT)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("FFFFFF")
    p.alignment = PP_ALIGN.CENTER

def _rail_card(slide, rail_left_in: float, rail_w_in: float, month_label: str, item: Dict) -> Tuple[float, float]:
    """
    Right-rail rounded 'card' with Month + ID + compact details.
    Returns (card_top_in, card_height_in).
    """
    if rail_w_in <= 0:
        return (0.0, 0.0)

    card_top_in = 1.35
    card_h_in   = 2.55
    card_pad_in = 0.35

    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(rail_left_in + 0.25),
        Inches(card_top_in),
        Inches(rail_w_in - 0.50),
        Inches(card_h_in),
    )
    shp.adjustments[0] = 0.18
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string("F8FAFC")
    shp.line.color.rgb = RGBColor.from_string("0B1220")
    shp.line.width = Pt(1.2)

    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(card_pad_in)
    tf.margin_right = Inches(card_pad_in)
    tf.margin_top = Inches(card_pad_in)
    tf.margin_bottom = Inches(card_pad_in)

    # Month
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = month_label or ""
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("6B21A8")

    # Roadmap ID
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    r = p.add_run()
    rid = str(item.get("roadmap_id") or item.get("roadmap") or "").strip()
    r.text = f"Roadmap ID: {rid}" if rid else ""
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor.from_string("111827")

    # Compact rows
    def add_row(label: str, value: str):
        if not value:
            return
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        r1 = p.add_run(); r1.text = f"{label} "; r1.font.bold = True
        r1.font.size = Pt(13); r1.font.color.rgb = RGBColor.from_string("111827")
        r2 = p.add_run(); r2.text = value
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor.from_string("374151")

    status_text = (item.get("status") or "").title()
    add_row("Status:",    status_text)
    add_row("GA:",        item.get("ga") or item.get("ga_end") or "")
    add_row("Platforms:", ", ".join(item.get("platforms", [])) if isinstance(item.get("platforms"), list) else (item.get("platforms") or ""))
    add_row("Clouds:",    ", ".join(item.get("clouds", []))    if isinstance(item.get("clouds"), list)    else (item.get("clouds") or ""))

    return card_top_in, card_h_in

def _normalize_audience(item: dict) -> Tuple[Optional[bool], Optional[bool]]:
    """
    Returns (end_users, admins) as booleans or None when unknown.
    Accepts dict/list/string forms from parsers.
    """
    aud = item.get("audience")
    eu = ad = None

    def yesish(v: str) -> bool:
        v = (v or "").strip().lower()
        return v in ("y","yes","true","1","✓","✔","☑") or "yes" in v

    if isinstance(aud, dict):
        if "end_users" in aud: eu = yesish(str(aud.get("end_users")))
        if "end user" in aud:  eu = yesish(str(aud.get("end user")))
        if "admins" in aud:    ad = yesish(str(aud.get("admins")))
        if "admin" in aud:     ad = yesish(str(aud.get("admin")))
    elif isinstance(aud, list):
        low = [str(x).lower() for x in aud]
        eu = True  if any("end user" in x or x == "users" for x in low) else eu
        ad = True  if any("admin" in x for x in low) else ad
    elif isinstance(aud, str):
        s = aud.lower()
        if any(k in s for k in ("end user","end-user","users")): eu = True
        if "admin" in s: ad = True

    return (eu, ad)

def _rail_audience_rows(slide, rail_left_in: float, rail_w_in: float, assets: Dict,
                        eu: Optional[bool], ad: Optional[bool],
                        card_top_in: float, card_h_in: float) -> None:
    """Header + two rows placed BELOW the rail card."""
    if rail_w_in <= 0:
        return

    left   = rail_left_in + 0.25
    row_w  = max(rail_w_in - 0.5, 1.4)
    row_h  = 0.55

    # Start just below the card
    start_top = card_top_in + card_h_in + 0.30

    # Header
    header_top = start_top
    hdr = slide.shapes.add_textbox(Inches(left), Inches(header_top),
                                   Inches(row_w), Inches(0.5)).text_frame
    p = hdr.paragraphs[0]
    r = p.add_run()
    r.text = "Target Audience:"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("CBD5E1")

    rows_top = header_top + 0.60
    rows = [
        ("End Users:", assets.get("icon_endusers"), eu),
        ("Admins:",    assets.get("icon_admins"),   ad),
    ]
    for idx, (label, icon, val) in enumerate(rows):
        top = rows_top + idx * (row_h + 0.20)
        # icon (optional; guard with exists)
        if icon and Path(icon).exists():
            slide.shapes.add_picture(icon, Inches(left), Inches(top),
                                     width=Inches(0.42), height=Inches(0.42))
        # label + value
        t = slide.shapes.add_textbox(Inches(left + 0.55), Inches(top),
                                     Inches(row_w - 0.55), Inches(row_h)).text_frame
        p = t.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{label} "
        r1.font.size = Pt(14); r1.font.bold = True
        r1.font.color.rgb = RGBColor.from_string("FFFFFF")

        r2 = p.add_run()
        val_txt = "Yes" if val is True else "?" if val is None else "No"
        r2.text = val_txt
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor.from_string("E5E7EB")

# ---- layout constants --------------------------------------------------------
LINE_HEIGHT_IN, MIN_TITLE_H, MAX_TITLE_H = 0.55, 1.10, 1.60
TITLE_FONT_PT = 44
BODY_FONT_PT  = 20
TITLE_TOP     = 0.80

def _estimate_title_height(text: str, width_in: float) -> float:
    """Crude estimate based on characters-per-line; clamps to MIN/MAX."""
    if not text:
        return MIN_TITLE_H
    # assume ~ 36 chars per 6.0in width, scale linearly
    cpp = 36 * (width_in / 6.0)
    lines = max(1, int((len(text) / max(1, cpp)) + 0.5))
    h = lines * (LINE_HEIGHT_IN * 0.92)
    return max(MIN_TITLE_H, min(MAX_TITLE_H, h))

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("-i", "--inputs", nargs="+", required=True, help="HTML inputs")
    ap.add_argument("-o", "--output", required=True, help="Output PPTX")
    ap.add_argument("--month", default="", help="Month label for slides")
    ap.add_argument("--template", default="", help="Optional PPTX template path")

    # assets (underscore style for run_build)
    ap.add_argument("--cover", default="", help="Cover background image")
    ap.add_argument("--agenda", default="", help="Agenda background image")
    ap.add_argument("--separator", default="", help="Separator background image")
    ap.add_argument("--conclusion", default="", help="Conclusion background image")
    ap.add_argument("--thankyou", default="", help="Thank-you background image")
    ap.add_argument("--brand_bg", default="", help="Brand background for item slides")
    ap.add_argument("--cover_title", default="M365 Technical Update Briefing", help="Cover title")
    ap.add_argument("--cover_dates", default="", help="Cover date text")
    ap.add_argument("--logo", default="", help="Primary logo")
    ap.add_argument("--logo2", default="", help="Secondary logo")

    # icons
    ap.add_argument("--icon_rocket", default="", help="Icon for GA/Rollout")
    ap.add_argument("--icon_preview", default="", help="Icon for Preview/In Dev")
    ap.add_argument("--icon_endusers", default="", help="Icon for End Users audience")
    ap.add_argument("--icon_admins", default="", help="Icon for Admins audience")

    ap.add_argument("--rail_width", type=float, default=None, help="Right rail width in inches (e.g., 3.5)")
    ap.add_argument("--debug_dump", default="", help="Write parsed items JSON here")

    args = ap.parse_args()

    # resolve assets
    def g(p: str) -> str:
        return p if p and Path(p).exists() else ""

    assets = {
        "cover":      g(args.cover),
        "agenda":     g(args.agenda),
        "separator":  g(args.separator),
        "conclusion": g(args.conclusion),
        "thankyou":   g(args.thankyou),
        "brand_bg":   g(args.brand_bg),
        "cover_title": args.cover_title,
        "cover_dates": args.cover_dates or args.month,
        "logo": g(args.logo),
        "logo2": g(args.logo2),
        # icons
        "icon_rocket":   g(args.icon_rocket),
        "icon_preview":  g(args.icon_preview),
        "icon_endusers": g(args.icon_endusers),
        "icon_admins":   g(args.icon_admins),
        # default palette (can be overridden by generate_deck.py/style)
        "product_palette": {
            "teams": "4F46E5",
            "sharepoint": "16A34A",
            "onedrive": "0EA5E9",
            "exchange": "F97316",
            "outlook": "2563EB",
            "default": "0F172A",
        },
    }

    build(
        inputs=args.inputs,
        output_path=args.output,
        month=args.month,
        assets=assets,
        template=args.template,
        rail_width=args.rail_width,
        conclusion_links=None,
        debug_dump=(args.debug_dump or None),
    )
'''
Path("/mnt/data/run_build.py").write_text(code, encoding="utf-8")
print("Wrote /mnt/data/run_build.py (full replacement).")
