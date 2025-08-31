# run_build.py
from pptx import Presentation
from layout import Layout
from parsers import parse_message_center_html, parse_roadmap_html
import os


def load_style(path: str) -> dict:
    with open(path, 'r') as f:
        return yaml.safe_load(f)


# run_build.py

from pptx import Presentation
from layout import Layout
from parsers import parse_message_center_html, parse_roadmap_html
import os
import logging

def build(
    inputs,
    output_path,
    month,
    assets,
    template=None,
    rail_width=3.5,
    conclusion_links=None
):
    print("Starting build process...")

    # Load or create presentation
    if template and os.path.exists(template):
        prs = Presentation(template)
        print(f"Loaded template: {template}")
    else:
        prs = Presentation()
        print("Created new presentation.")

    # Parse all HTML inputs
    all_items = []
    for file_path in inputs:
        if not os.path.exists(file_path):
            print(f"Input file not found: {file_path}")
            continue
        filename_lower = file_path.lower()
        if "messagecenter" in filename_lower:
            parsed = parse_message_center_html(file_path, month)
        elif "roadmap" in filename_lower:
            parsed = parse_roadmap_html(file_path, month)
        else:
            print(f"Unrecognized input type: {file_path}")
            parsed = []
        print(f"Parsed {len(parsed)} items from {file_path}")
        all_items.extend(parsed)

    # Deduplicate items based on roadmap_id, title, or URL
    seen = set()
    unique_items = []
    for item in all_items:
        key = (item.get("roadmap_id") or item.get("title") or item.get("url") or "").lower()
        if key and key not in seen:
            seen.add(key)
            unique_items.append(item)
    print(f"Total unique items: {len(unique_items)}")

    # Sort items (optional)
    unique_items.sort(key=lambda i: (i.get("products", [""])[0], i.get("title", "")))

    # Initialize layout with style
    style_cfg = load_style("style_template.yaml")
    layout = Layout(style_cfg)

    # Add cover slide
    cover_title = assets.get("cover_title", "My Quarterly Update")
    cover_dates = assets.get("cover_dates", month)
    layout.add_cover_slide(prs, assets, {"title": cover_title, "dates": cover_dates})

    # Group items by product (for sections)
    grouped = {}
    order = []
    for item in unique_items:
        products = item.get("products", [])
        product_name = products[0] if products else "General"
        if product_name not in grouped:
            grouped[product_name] = []
            order.append(product_name)
        grouped[product_name].append(item)

    # Generate section slides per product
    for product in order:
        layout.add_separator_slide(prs, assets, f"{product} updates")
        for item in grouped[product]:
            layout.add_item_slide(prs, item, month, assets)

    # Add conclusion slide
    links = conclusion_links or [
        ("Security", "https://www.microsoft.com/security"),
        ("Azure", "https://azure.microsoft.com/"),
        ("Docs", "https://learn.microsoft.com/")
    ]
    layout.add_conclusion_slide(prs, assets, links)

    # Add thank you slide
    layout.add_thankyou_slide(prs, assets)

    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    """
    Generate PPTX from input HTML files.
    """
    print("Starting deck build...")
    # Aggregate parsed items
    all_items = []
    for file_path in inputs:
        ext = os.path.splitext(file_path)[1].lower()
        if "messagecenter" in file_path.lower():
            parsed = parse_message_center_html(file_path, month)
        elif "roadmap" in file_path.lower():
            parsed = parse_roadmap_html(file_path, month)
        else:
            parsed = []  # For unknown types or extend as needed
        all_items.extend(parsed)

    # Deduplicate
    seen = set()
    unique_items = []
    for item in all_items:
        key = (item.get("roadmap_id") or item.get("title") or item.get("url") or "").lower()
        if key and key not in seen:
            seen.add(key)
            unique_items.append(item)

    # Sort for presentation flow (by product, title)
    unique_items.sort(key=lambda i: (i.get("products", []), i.get("title", "")))

    # Create presentation
    if template and os.path.exists(template):
        prs = Presentation(template)
    else:
        prs = Presentation()

    # Initialize Layout with style
    style_cfg = load_style("style_template.yaml")
    layout = Layout(style_cfg)

    # Add cover slide
    cover_title = assets.get("cover_title", "My Quarterly Update")
    cover_dates = assets.get("cover_dates", month)
    layout.add_cover_slide(prs, assets, {"title": cover_title, "dates": cover_dates})

    # Group items by product
    grouped = {}
    order = []
    for item in unique_items:
        product = item.get("products", ["General"])[0] if item.get("products") else "General"
        if product not in grouped:
            grouped[product] = []
            order.append(product)
        grouped[product].append(item)

    # Generate product section slides
    for product in order:
        layout.add_separator_slide(prs, assets, f"{product} updates")
        for item in grouped[product]:
            layout.add_item_slide(prs, item, month, assets)

    # Add conclusion slide
    links = conclusion_links or [
        ("Security", "https://www.microsoft.com/security"),
        ("Azure", "https://azure.microsoft.com/"),
        ("Docs", "https://learn.microsoft.com/")
    ]
    layout.add_conclusion_slide(prs, assets, links)

    # Add thank you
    layout.add_thankyou_slide(prs, assets)

    # Save presentation
    prs.save(output_path)
    print(f"Deck saved to {output_path}")