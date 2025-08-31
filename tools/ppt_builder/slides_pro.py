# slides_pro.py - closer to Sample deck look
import os
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE

def inches_to_emu(x): return int(x * 914400)

def _hex(color):
    s = str(color or "FFFFFF").strip()
    if s.startswith("#"): s = s[1:]
    if len(s) == 3: s = "".join([c*2 for c in s])
    if len(s) != 6: s = "FFFFFF"
    return s.upper()

def _font_run(p, text, size, bold=False, color="FFFFFF", name=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    if name: r.font.name = name
    r.font.color.rgb = RGBColor.from_string(_hex(color))
    return r

def add_full_bg(slide, image_path):
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, width=inches_to_emu(10), height=inches_to_emu(7.5))
    else:
        # fallback: a subtle dark fill rectangle
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inches_to_emu(10), inches_to_emu(7.5))
        fill = shp.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(22,28,36)  # #161C24
        shp.line.fill.background()

def add_title(slide, text, top=0.9, size=52, align=PP_ALIGN.LEFT, color="#FFFFFF", name="Segoe UI Semibold"):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(1.2))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    _font_run(p, text or "", size, bold=True, color=color, name=name)

def add_text(slide, text, left=0.6, top=2.0, width=8.8, height=3.4, size=20, color="#FFFFFF", name="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    _font_run(p, text or "", size, bold=False, color=color, name=name)

def add_chip(slide, text, left, top, fill="#2E7D32", text_color="#FFFFFF"):
    """Rounded chip for statuses like 'Preview', 'GA' etc."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches_to_emu(left), inches_to_emu(top), inches_to_emu(1.6), inches_to_emu(0.45))
    shp.adjustments[0] = 0.5  # more rounded
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(_hex(fill))
    shp.line.fill.background()
    tb = shp.text_frame
    tb.clear(); tb.word_wrap = True
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _font_run(p, text, 14, bold=True, color=text_color, name="Segoe UI Semibold")

def cover(prs, assets, title, dates, logo1=None, logo2=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("cover"))
    add_title(slide, title or "M365 Technical Update Briefing", top=2.1, size=54, align=PP_ALIGN.CENTER)
    if dates:
        add_text(slide, dates, left=0.6, top=3.8, width=8.8, height=0.8, size=30, color="#E6E8EF", name="Segoe UI",)
    # logos
    if logo1 and os.path.exists(logo1):
        slide.shapes.add_picture(logo1, inches_to_emu(0.4), inches_to_emu(6.6), height=inches_to_emu(0.6))
    if logo2 and os.path.exists(logo2):
        slide.shapes.add_picture(logo2, inches_to_emu(8.0), inches_to_emu(6.6), height=inches_to_emu(0.6))

def agenda(prs, assets, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("agenda"))
    add_title(slide, "Agenda", top=0.9)
    top = 2.1
    for line in lines:
        add_text(slide, f"• {line}", left=0.9, top=top, width=8.2, height=0.5, size=26)
        top += 0.55

def legend(prs, assets):
    """Legend slide showing Target Audience, Cadence, Feature Status similar to sample."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("brand_bg") or assets.get("separator"))
    add_title(slide, "Legend")
    # three columns
    col_w = 2.8
    # Target Audience
    add_text(slide, "Target Audience:", left=0.6, top=2.0, width=col_w, height=0.5, size=22)
    add_text(slide, "Admins\nEnd-users\nDevelopers\nSecurity\nCompliance", left=0.6, top=2.6, width=col_w, height=2.0, size=18, color="#E6E8EF")
    # Cadence
    add_text(slide, "Cadence:", left=3.6, top=2.0, width=col_w, height=0.5, size=22)
    add_text(slide, "Monthly\nQuarterly\nYearly", left=3.6, top=2.6, width=col_w, height=2.0, size=18, color="#E6E8EF")
    # Feature Status
    add_text(slide, "Feature Status:", left=6.6, top=2.0, width=col_w, height=0.5, size=22)
    add_chip(slide, "Preview", left=6.6, top=2.7, fill="#1976D2")
    add_chip(slide, "GA", left=8.0, top=2.7, fill="#2E7D32")
    add_chip(slide, "Rollout", left=6.6, top=3.3, fill="#8E24AA")
    add_chip(slide, "Deprecated", left=8.0, top=3.3, fill="#B71C1C")

def separator(prs, assets, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("separator"))
    add_title(slide, title, top=3.0, size=56, align=PP_ALIGN.CENTER)

def item(prs, assets, item_dict, month):
    """Two-row layout: title + body, with chips rail at bottom"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("brand_bg") or assets.get("cover"))
    title = item_dict.get("title","")
    body = item_dict.get("summary") or item_dict.get("description") or ""
    add_title(slide, title, top=0.6, size=38)
    add_text(slide, body, left=0.6, top=2.0, width=8.8, height=3.6, size=20, color="#FFFFFF")
    # meta chips
    status = (item_dict.get("status") or "").lower()
    if "preview" in status:
        add_chip(slide, "Preview", left=0.6, top=6.1, fill="#1976D2")
    if "ga" in status or "general availability" in status:
        add_chip(slide, "GA", left=2.3, top=6.1, fill="#2E7D32")
    if "rolling" in status or "rollout" in status:
        add_chip(slide, "Rollout", left=3.7, top=6.1, fill="#8E24AA")
    rid = item_dict.get("roadmap_id")
    if rid:
        add_chip(slide, f"ID {rid}", left=5.3, top=6.1, fill="#455A64")
    if month:
        add_chip(slide, month, left=7.2, top=6.1, fill="#263238")

def conclusion(prs, assets, links):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("conclusion"))
    add_title(slide, "Final Thoughts", top=0.9)
    top = 2.3
    for text, url in links:
        add_text(slide, f"{text}: {url}", left=0.6, top=top, width=8.8, height=0.6, size=22, color="#E6E8EF")
        top += 0.6

def thankyou(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("thankyou"))

