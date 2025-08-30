from __future__ import annotations

import argparse
import os
import re
import time
from dataclasses import dataclass
from datetime import datetime
from typing import Any

from bs4 import BeautifulSoup
from bs4.element import NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# new import
from parsers import parse_message_center_html, parse_roadmap_html
from dataclasses import fields as _dc_fields


# ----------------------------
# Constants / theme
# ----------------------------
EMU_PER_INCH = 914400

GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_PURPLE = RGBColor(40, 17, 63)  # rail
MID_PURPLE = RGBColor(70, 35, 100)
LIGHT_PURPLE = RGBColor(95, 60, 135)

DEFAULT_RAIL_WIDTH_IN = 3.5
PAGE_MARGIN_IN = 0.6


# ----------------------------
# Utils
# ----------------------------
EMU_PER_INCH = 914400  # python-pptx constant

def emu_to_inches(x: int | float | None) -> float:
    """Convert EMU to inches; None -> 0.0."""
    if x is None:
        return 0.0
    return float(x) / EMU_PER_INCH

def inches_to_emu(x: float | int | None) -> int:
    """Convert inches to EMU; None -> 0.0."""
    if x is None:
        return 0
    return int(round(float(x) * EMU_PER_INCH))


# --- stdlib
import os, sys, re, traceback
from typing import Any, List, Dict, Optional

# --- third-party
from bs4 import BeautifulSoup

# --- bs4-safe helpers
def _txt(x: Any) -> str:
    try:
        return x.get_text(strip=True)  # type: ignore[attr-defined]
    except Exception:
        return "" if x is None else str(x)

def _attr(x: Any, name: str) -> str:
    try:
        v = x.get(name)  # type: ignore[call-arg, attr-defined]
        if v is None:
            return ""
        if isinstance(v, str):
            return v
        # bs4 can return list-like attrs
        try:
            return " ".join(v)
        except Exception:
            return str(v)
    except Exception:
        return ""

def first(el: Any, css: str) -> Any:
    try:
        return el.select_one(css)  # type: ignore[attr-defined]
    except Exception:
        return None

def all_of(el: Any, css: str) -> List[Any]:
    try:
        res = el.select(css)  # type: ignore[attr-defined]
        return list(res) if res else []
    except Exception:
        return []

   

def _to_item(d: dict[str, Any]) -> Item:
    valid = {f.name for f in _dc_fields(Item)}
    data: dict[str, Any] = {k: d.get(k) for k in valid}

    # Normalize core string fields
    for k in ("title", "summary", "url", "roadmap_id", "status", "month"):
        v = data.get(k)
        data[k] = "" if v is None else str(v)

    # Normalize list-ish fields to list[str]
    list_fields = ("products", "platforms", "audience", "phases", "clouds")
    for k in list_fields:
        v = data.get(k)
        if v is None:
            data[k] = []
        elif isinstance(v, (list, tuple, set)):
            data[k] = [str(x) for x in v if x is not None]
        else:
            # Single value -> wrap
            data[k] = [str(v)]

    # Optional date-ish fields -> strings
    for k in ("created", "modified", "ga"):
        v = data.get(k)
        if v is None:
            data[k] = ""
        else:
            data[k] = str(v)

    return Item(**data)




def _to_text(val: Any) -> str:
    """Turn bs4 values (str | list[str] | Tag | NavigableString | None) into text."""
    if val is None:
        return ""
    if isinstance(val, NavigableString,):
        return str(val)
    if isinstance(val, Tag):
        return val.get_text(" ", strip=True)
    from collections.abc import Sequence
    # treat true sequences but not text-like
    if isinstance(val, Sequence) and not isinstance(val, (str, bytes, bytearray)):
         return ", ".join(_to_text(v) for v in val if v is not None)
    return str(val)

# after: from pptx.util import Inches, Pt
from typing import Optional

def _inches(x: Optional[float]) -> int:
    """Return EMU from inches, accepting None by coercing to 0.0."""
    return _inches(0.0 if x is None else float(x))

def _pt(x: Optional[float]) -> int:
    """Return EMU from points, accepting None by coercing to 0.0."""
    return Pt(0.0 if x is None else float(x))


def smart_split_product_title(title: str) -> tuple[str, str]:
    """
    If title is like 'Product: Feature name', return ('Product', 'Feature name').
    Otherwise return ('', title).
    """
    if not title:
        return "", ""
    m = re.match(r"^\s*([^:]{2,50})\s*:\s*(.+)$", title)
    if m:
        prod = m.group(1).strip()
        rest = m.group(2).strip()
        return prod, rest
    return "", title.strip()


def pick_status_icon_key(status_text: str) -> str:
    """
    Choose 'rocket' for GA/Launched/Rolling out; 'magnifier' for Preview/In development/Planned.
    Default to 'magnifier' when uncertain.
    """
    s = (status_text or "").lower()
    if any(k in s for k in ["launched", "rolling out", "rolled out", "general availability", "ga"]):
        return "rocket"
    if any(k in s for k in ["preview", "in development", "planned", "beta"]):
        return "magnifier"
    return "magnifier"


def path_if_exists(p: str | None) -> str | None:
    if p and os.path.exists(p):
        return p
    return None


def add_picture_safe(
    slide,
    image_path: str | None,
    left_in: float,
    top_in: float,
    width_in: float | None = None,
    height_in: float | None = None,
):
    """Add picture if file is present. width_in/height_in are optional (Inches)."""
    if not image_path or not os.path.exists(image_path):
        return None
    left = _inches(left_in)
    top = _inches(top_in)
    if width_in is None and height_in is None:
        return slide.shapes.add_picture(image_path, left, top)
    if width_in is not None and height_in is None:
        return slide.shapes.add_picture(image_path, left, top, width=_inches(width_in))
    if width_in is None and height_in is not None:
        return slide.shapes.add_picture(image_path, left, top, height=_inches(height_in))
    return slide.shapes.add_picture(
        image_path, left, top, width=_inches(width_in), height=_inches(height_in)
    )



def safe_find(node: Tag | None, name, **kwargs) -> Tag | None:
    if not _is_tag(node):
        return None
    t = node.find(name, **kwargs)
    return t if _is_tag(t) else None

def safe_find_all(node: Tag | None, name, **kwargs) -> list[Tag]:
    if not _is_tag(node):
        return []
    return [t for t in node.find_all(name, **kwargs) if _is_tag(t)]


def add_title_box(
    slide,
    text: str,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    font_size_pt: int = 60,
    bold: bool = True,
    color: RGBColor = GOLD,
    align=PP_ALIGN.LEFT,
):
    """Title that wraps & shrinks to fit."""
    box = slide.shapes.add_textbox(
        _inches(left_in), _inches(top_in), _inches(width_in), _inches(height_in)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = _inches(0.05)
    tf.margin_right = _inches(0.05)
    tf.margin_top = _inches(0.02)
    tf.margin_bottom = _inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text or ""
    f = r.font
    f.size = Pt(font_size_pt)
    f.bold = bold
    f.color.rgb = color
    return box


def add_text_box(
    slide,
    text: str,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    size_pt: int = 18,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(
        _inches(left_in), _inches(top_in), _inches(width_in), _inches(height_in)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text or ""
    f = r.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = color
    return box


def add_full_slide_picture(slide, prs, image_path: str | None):
    """Stretch a background image to full slide; ignore if missing."""
    if not image_path or not os.path.exists(image_path):
        return
    sw = emu_to_inches(prs.slide_width)
    sh = emu_to_inches(prs.slide_height)
    add_picture_safe(slide, image_path, left_in=0.0, top_in=0.0, width_in=sw, height_in=sh)


def draw_side_rail(
    slide, prs, rail_left_in: float, rail_width_in: float, color: RGBColor = DARK_PURPLE
):
    """Right-side vertical rail rectangle, sized using the Presentation dimensions."""
    slide_h_in = emu_to_inches(prs.slide_height)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        _inches(rail_left_in),
        _inches(0.0),
        _inches(rail_width_in),
        _inches(slide_h_in),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bubble(slide, text: str, left_in: float, top_in: float, width_in: float, height_in: float):
    """Rounded rectangle bubble for product/technology."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _inches(left_in),
        _inches(top_in),
        _inches(width_in),
        _inches(height_in),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = MID_PURPLE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(16)
    f.bold = True
    f.color.rgb = WHITE
    return shape


def add_notes(slide, text: str):
    if not text:
        return
    notes = slide.notes_slide.notes_text_frame
    if notes.paragraphs and notes.paragraphs[0].text:
        notes.text += "\n\n" + text
    else:
        notes.text = text


def _safe_save(prs, output_path: str, tries: int = 4, delay_sec: float = 0.6) -> str:
    base, ext = os.path.splitext(output_path)
    for i in range(tries):
        candidate = (
            output_path if i == 0 else f"{base}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{i}{ext}"
        )
        try:
            prs.save(candidate)
            return candidate
        except PermissionError:
            time.sleep(delay_sec)
            continue
    prs.save(output_path)
    return output_path


# ----------------------------
# Data model
# ----------------------------

# ---------- CORE MODEL & SAFE HELPERS ----------


from dataclasses import dataclass, field, asdict
from typing import List, Optional, Iterable, Any, Sequence
import os, csv, re, datetime as _dt

from bs4 import BeautifulSoup  # keep BS import simple
from bs4.element import Tag as Bs4Tag, NavigableString as Bs4Nav, PageElement as Bs4El, ResultSet as Bs4RS

# Colors / constants (optional)
PHASE_TO_STATUS_MAP = {
    "in development": "In development",
    "rolling out": "Rolling out",
    "launched": "Launched",
    "preview": "Preview",
    "cancelled": "Cancelled",
    "delayed": "Delayed",
}

# --- Data model used by slide builder ---
@dataclass
class Item:
    # required-ish
    title: str = ""
    summary: str = ""          # short synopsis
    description: str = ""      # longer body if present
    roadmap_id: str = ""       # "MC#####" or similar / roadmap numeric id
    url: str = ""              # canonical link
    month: str = ""            # e.g., "September 2025"

    # categorization
    product: str = ""          # primary product string
    products: List[str] = field(default_factory=list)
    platforms: List[str] = field(default_factory=list)  # Win/Mac/iOS/Web, etc.
    audience: List[str] = field(default_factory=list)   # Admins/Users/etc.
    clouds: List[str] = field(default_factory=list)     # GCC/GCC-H/DoD/Commercial, etc.

    # lifecycle
    status: str = ""           # In development / Rolling out / Preview / Launched / …
    phases: str = ""           # raw phases text if source uses it
    created: str = ""          # dates as plain strings is fine for slides
    modified: str = ""
    ga: str = ""               # GA/available date if present

def _is_tag(x: Any) -> bool:
    return isinstance(x, Bs4Tag)

def _as_tag(x: Any) -> Optional[Bs4Tag]:
    return x if isinstance(x, Bs4Tag) else None

def safe_text(x: Any) -> str:
    """Get visible text from Tag/NavigableString/other safely."""
    if x is None:
        return ""
    if isinstance(x, (Bs4Tag, Bs4Nav)):
        try:
            return x.get_text(" ", strip=True)
        except Exception:
            return str(x).strip()
    return str(x).strip()

def attr_str(tag: Optional[Bs4Tag], key: str) -> str:
    """Return attribute as string; handles list-valued attributes."""
    if not _is_tag(tag):
        return ""
    val = tag.get(key)  # type: ignore[attr-defined]
    if val is None:
        return ""
    if isinstance(val, list):
        return " ".join([str(v) for v in val])
    return str(val)

def find_one(scope: Bs4Tag, name: Any = None, **kwargs: Any) -> Optional[Bs4Tag]:
    """BeautifulSoup find() but always returns Tag or None."""
    try:
        return _as_tag(scope.find(name, **kwargs))  # type: ignore[no-untyped-call]
    except Exception:
        return None

def find_all_tags(scope: Bs4Tag, name: Any = None, **kwargs: Any) -> List[Bs4Tag]:
    """BeautifulSoup find_all() but filtered to Tag instances."""
    try:
        return [t for t in scope.find_all(name, **kwargs) if _is_tag(t)]  # type: ignore[no-untyped-call]
    except Exception:
        return []

def first_or_none(seq: Iterable[Any]) -> Optional[Any]:
    for x in seq:
        return x
    return None

def _clean(s: Optional[str]) -> str:
    if not s:
        return ""
    # collapse whitespace, strip control chars
    s = re.sub(r"\s+", " ", s).strip()
    # prevent weird NBSPs etc.
    return s.replace("\u00A0", " ")

def clamp(s: str, max_len: int) -> str:
    if len(s) <= max_len:
        return s
    return s[: max(0, max_len - 1)].rstrip() + "…"

def map_phase_to_status(phase_text: str) -> str:
    k = _clean(phase_text).lower()
    for key, val in PHASE_TO_STATUS_MAP.items():
        if key in k:
            return val
    return _clean(phase_text)  # fallback: return original string

def normalize_item(i: Item, month_fallback: str = "") -> Item:
    """Ensure required fields, map synonyms, fill lists, trim for layout."""
    i.title = clamp(_clean(i.title), 140)
    i.summary = _clean(i.summary)
    i.description = _clean(i.description)
    i.roadmap_id = _clean(i.roadmap_id)
    i.url = _clean(i.url)

    # month
    if not _clean(i.month):
        i.month = _clean(month_fallback)

    # status from phases if missing
    if not _clean(i.status) and _clean(i.phases):
        i.status = map_phase_to_status(i.phases)

    # products list should include primary product
    if _clean(i.product) and _clean(" ".join(i.products)) == "":
        i.products = [i.product]
    elif _clean(i.product) and i.product not in i.products:
        i.products.insert(0, i.product)

    # audience normalize (accept string or list upstream)
    i.audience = [a for a in (i.audience or []) if _clean(a)]
    i.platforms = [p for p in (i.platforms or []) if _clean(p)]
    i.clouds = [c for c in (i.clouds or []) if _clean(c)]

    return i

# ---------- REPORTING HELPERS ----------
def write_parse_report(items: List[Item], out_dir: str) -> str:
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, "last_parse_report.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write(f"Parse report ({_dt.datetime.now().isoformat(timespec='seconds')}):\n")
        f.write(f"Total items: {len(items)}\n\n")
        for idx, it in enumerate(items[:10], start=1):  # show first 10
            f.write(f"[{idx}] {it.title}\n")
            f.write(f"    ID: {it.roadmap_id} | Status: {it.status} | Month: {it.month}\n")
            f.write(f"    Product: {it.product} | Products: {', '.join(it.products)}\n")
            f.write(f"    Platforms: {', '.join(it.platforms)} | Audience: {', '.join(it.audience)}\n")
            f.write(f"    URL: {it.url}\n\n")
    return path

def write_items_csv(items: List[Item], path: str) -> str:
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
    cols = [
        "title","summary","description","roadmap_id","url","month",
        "product","products","platforms","audience","clouds",
        "status","phases","created","modified","ga",
    ]
    with open(path, "w", newline="", encoding="utf-8") as fp:
        w = csv.writer(fp)
        w.writerow(cols)
        for it in items:
            row = asdict(it)
            row["products"] = ";".join(row.get("products", []) or [])
            row["platforms"] = ";".join(row.get("platforms", []) or [])
            row["audience"] = ";".join(row.get("audience", []) or [])
            row["clouds"] = ";".join(row.get("clouds", []) or [])
            w.writerow([row.get(k, "") for k in cols])
    return path



# ----------------------------
# HTML parsing
# ----------------------------
STATUS_WORDS = [
    "Launched",
    "Rolling out",
    "General Availability",
    "GA",
    "Preview",
    "Public Preview",
    "Private Preview",
    "In development",
    "Planned",
    "Beta",
]

AUDIENCE_WORDS = ["Admin", "Administrator", "IT Admin", "End user", "User", "Developer"]


def _mk_item(**kwargs) -> Item:
    filtered = {k: v for k, v in kwargs.items() if k in _allowed_fields}
    return Item(**filtered)  # type: ignore[call-arg]

def _classes(node) -> list[str]:
    c = None
    try:
        c = node.get("class")
    except Exception:
        pass
    if not c:
        return []
    if isinstance(c, str):
        return c.split()
    if isinstance(c, (list, tuple, set)):
        return [str(x) for x in c]
    return [str(c)]

def _find_url(card) -> str:
    for a in safe_find_all(card, "a", href=True):
        href = _attr(a, "href") or ""
        if re.search(r"(featureid=|\broadmap\b|\bmicrosoft-365-roadmap\b)", href, re.I):
            return href
    a0 = safe_find(card, "a", href=True)
    return (_attr(a0, "href") or "") if a0 else ""

def _find_title(card) -> str:
    # Prefer headings
    for h in ("h1", "h2", "h3", "h4"):
        htag = safe_find(card, h)
        if htag:
            t = _txt(htag)
            if t:
                return t
    # Next, any element with role=heading or strong/b
    role_h = safe_find(card, attrs={"role": "heading"})
    if role_h:
        t = _txt(role_h)
        if t:
            return t
    sb = safe_find(card, "strong") or safe_find(card, "b")
    if sb:
        t = _txt(sb)
        if t:
            return t
    # Fallback to first link text or the whole card text
    a = safe_find(card, "a")
    if a:
        t = _txt(a)
        if t:
            return t
    return _txt(card)

def _longest_para(card) -> str:
    paras = [p for p in safe_find_all(card, ["p", "div", "span"]) if _txt(p)]
    if not paras:
        return _txt(card)
    paras.sort(key=lambda p: len(_txt(p)), reverse=True)
    return _txt(paras[0])

def _find_label_value(card, label: str) -> str:
    """
    Find a value following a label like 'Products', 'Platforms', etc.
    Tries multiple structures:
        <span>Products</span><span>Excel;Teams</span>
        <div class='field'><label>Products</label><div>Excel</div></div>
    """
    lab_lower = label.lower()
    # Any tag whose text equals the label
    for tag in safe_find_all(card, True):
        txt = (_txt(tag) or "").strip().lower()
        if txt == lab_lower or txt == f"{lab_lower}:":
            # try next siblings or parent pattern
            sib = getattr(tag, "find_next_sibling", None)
            if callable(sib):
                nxt = tag.find_next_sibling()
                if nxt:
                    val = _txt(nxt)
                    if val:
                        return val
            parent = getattr(tag, "parent", None)
            if parent:
                # look for a value element within same group
                cand = safe_find(parent, ["span", "div"])
                if cand and cand is not tag:
                    val = _txt(cand)
                    if val and val.lower() != lab_lower:
                        return val
    # Regex scan fallback like "Products: Excel, Teams"
    text = _txt(card) or ""
    m = re.search(rf"\b{re.escape(label)}\s*:\s*(.+)", text, re.I)
    if m:
        return _clean(m.group(1))
    return ""

def _find_feature_id(card) -> str:
    # From URL first
    url = _find_url(card)
    m = re.search(r"[?&]featureid=(\d+)", url, re.I)
    if m:
        return m.group(1)

    # From labeled text "Feature ID:"
    txt = _txt(card) or ""
    m2 = re.search(r"\bFeature\s*ID\s*:\s*(\d+)", txt, re.I)
    if m2:
        return m2.group(1)

    # From any visible "RM" code-ish token like "RM123456"
    m3 = re.search(r"\bRM(\d{4,})\b", txt, re.I)
    if m3:
        return m3.group(1)

    return ""




# ----------------------------
# Slide builders
# ----------------------------
def add_bg(slide, prs, image_path: str | None):
    add_full_slide_picture(slide, prs, image_path)


def add_cover_slide(
    prs,
    assets: dict,
    cover_title: str | None,
    cover_dates: str | None,
    logo1: str | None,
    logo2: str | None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, prs, assets.get("cover"))



    # Dates
    add_text_box(
        slide,
        cover_dates or "",
        left_in=left,
        top_in=top + height + 0.2,
        width_in=width,
        height_in=0.6,
        size_pt=28,
        bold=False,
        color=WHITE,
    )

    # Logos bottom-left/right
    add_picture_safe(slide, path_if_exists(logo1), left_in=0.4, top_in=6.6, height_in=0.6)
    add_picture_safe(slide, path_if_exists(logo2), left_in=sw_in - 2.2, top_in=6.5, width_in=1.8)


def add_agenda_slide(prs, assets: dict, agenda_lines: list[str] | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, assets.get("agenda"))
    # Title
    sw_in = emu_to_inches(prs.slide_width)
    add_title_box(
        slide,
        "Agenda",
        left_in=PAGE_MARGIN_IN,
        top_in=0.9,
        width_in=sw_in - 2 * PAGE_MARGIN_IN,
        height_in=1.2,
        font_size_pt=52,
        color=GOLD,
    )
    # Bullets
    if not agenda_lines:
        agenda_lines = ["Overview", "Key updates by product", "Timeline & rollout status", "Q&A"]
    top = 2.4
    for line in agenda_lines:
        add_text_box(
            slide,
            f"• {line}",
            left_in=PAGE_MARGIN_IN,
            top_in=top,
            width_in=sw_in - 2 * PAGE_MARGIN_IN,
            height_in=0.5,
            size_pt=26,
            color=WHITE,
        )
        top += 0.6


def add_separator_slide(prs, assets: dict, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, assets.get("separator"))
    sw_in = emu_to_inches(prs.slide_width)
    add_title_box(
        slide,
        title,
        left_in=PAGE_MARGIN_IN,
        top_in=3.2,
        width_in=sw_in - 2 * PAGE_MARGIN_IN,
        height_in=1.5,
        font_size_pt=56,
        color=GOLD,
    )


def add_conclusion_slide(prs, assets: dict, links: list[tuple[str, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, assets.get("conclusion"))
    sw_in = emu_to_inches(prs.slide_width)
    add_title_box(
        slide,
        "Final Thoughts",
        left_in=PAGE_MARGIN_IN,
        top_in=0.9,
        width_in=sw_in - 2 * PAGE_MARGIN_IN,
        height_in=1.2,
        font_size_pt=52,
        color=GOLD,
    )
    top = 2.4
    for text, url in links:
        add_text_box(
            slide,
            f"{text}: {url}",
            left_in=PAGE_MARGIN_IN,
            top_in=top,
            width_in=sw_in - 2 * PAGE_MARGIN_IN,
            height_in=0.5,
            size_pt=22,
            color=WHITE,
        )
        top += 0.6


def add_thankyou_slide(prs, assets: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, assets.get("thankyou"))
    # nothing else required; background handles the design


def add_item_slide(
    prs: Presentation,
    it,
    month_str: str = "",
    assets: dict | None = None,
    rail_left_in: float = 3.5,
    rail_width_in: float = 3.5,
) -> None:
    """
    Add a single roadmap/message-center item slide.

    Expects `it` to have (best-effort):
      - title: str
      - summary or description: str
      - roadmap_id or rid: str
      - url: str
      - status: str
      - products: List[str]
      - platforms: List[str]
      - audience: str
    """
    assets = assets or {}

    # --- helpers -------------------------------------------------------------
    def _emu_to_inches(emu: int) -> float:
        # 1 inch = 914400 EMUs
        return float(emu) / 914400.0

    def _hex_to_rgb(s: str | None, fallback=(59, 46, 90)) -> RGBColor:
        # default: a dark purple-ish
        if not s:
            r, g, b = fallback
            return RGBColor(r, g, b)
        ss = s.strip().lstrip("#")
        if len(ss) == 6:
            r = int(ss[0:2], 16)
            g = int(ss[2:4], 16)
            b = int(ss[4:6], 16)
            return RGBColor(r, g, b)
        r, g, b = fallback
        return RGBColor(r, g, b)

    def _first_nonempty(*vals: str) -> str:
        for v in vals:
            if v:
                x = str(v).strip()
                if x:
                    return x
        return ""

    def _as_csv(x) -> str:
        if x is None:
            return ""
        if isinstance(x, (list, tuple, set)):
            return ", ".join(str(i) for i in x if i)
        return str(x)

    def _pick_title_size(t: str) -> int:
        L = len(t)
        # simple heuristic to avoid overflow without text measurement
        if L <= 48:
            return 36
        if L <= 72:
            return 32
        if L <= 100:
            return 28
        if L <= 140:
            return 24
        return 22

    # --- pull data from item -------------------------------------------------
    title = _first_nonempty(getattr(it, "title", ""), getattr(it, "name", ""), "Untitled")
    summary = _first_nonempty(getattr(it, "summary", ""), getattr(it, "description", ""))
    rid = _first_nonempty(getattr(it, "roadmap_id", ""), getattr(it, "rid", ""))
    url = str(getattr(it, "url", "") or "")
    status = str(getattr(it, "status", "") or "")
    products = getattr(it, "products", None) or []
    platforms = getattr(it, "platforms", None) or []
    audience = str(getattr(it, "audience", "") or "")

    # --- slide geometry ------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw_in = _emu_to_inches(prs.slide_width)
    sh_in = _emu_to_inches(prs.slide_height)

    # Clamp/normalize rail metrics
    rail_left = float(rail_left_in or 0.0)
    rail_width = max(0.5, float(rail_width_in or 3.5))
    rail_left = max(0.0, min(rail_left, max(0.0, sw_in - rail_width)))

    content_left = rail_left + rail_width + 0.25
    content_width = max(2.5, sw_in - content_left - 0.4)

    # --- side rail rectangle -------------------------------------------------
    rail_color = _hex_to_rgb(assets.get("rail_color") or "#3B2E5A")
    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(rail_left),
        Inches(0.0),
        Inches(rail_width),
        Inches(sh_in),
    )
    fill = rail.fill
    fill.solid()
    fill.fore_color.rgb = rail_color
    rail.line.fill.background()  # no border

    # --- title ---------------------------------------------------------------
    title_size = _pick_title_size(title)
    tx = slide.shapes.add_textbox(
        Inches(content_left),
        Inches(0.6),
        Inches(content_width),
        Inches(1.8),
    )
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

    # --- summary/body --------------------------------------------------------
    if summary:
        body = slide.shapes.add_textbox(
            Inches(content_left),
            Inches(2.1),
            Inches(content_width),
            Inches(3.3),
        )
        btf = body.text_frame
        btf.clear()
        p = btf.paragraphs[0]
        p.text = summary
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(32, 32, 32)
        p.alignment = PP_ALIGN.LEFT

    # --- meta on the rail ----------------------------------------------------
    meta_left = rail_left + 0.35
    meta_width = rail_width - 0.7
    meta_top = 0.6
    line_gap = 0.35

    def _add_meta(label: str, value: str):
        nonlocal meta_top
        if not value:
            return
        box = slide.shapes.add_textbox(
            Inches(meta_left),
            Inches(meta_top),
            Inches(meta_width),
            Inches(0.35),
        )
        mtf = box.text_frame
        mtf.clear()
        # label
        p1 = mtf.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        # value
        p2 = mtf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(230, 230, 230)
        p2.space_before = Pt(1)
        meta_top += line_gap

    _add_meta("Feature ID", rid)
    _add_meta("Status", status)
    _add_meta("Products", _as_csv(products))
    _add_meta("Platforms", _as_csv(platforms))
    _add_meta("Audience", audience)
    _add_meta("Month", month_str or "")

    if url:
        # URL box (slightly more space)
        box = slide.shapes.add_textbox(
            Inches(meta_left),
            Inches(meta_top),
            Inches(meta_width),
            Inches(0.45),
        )
        mtf = box.text_frame
        mtf.clear()
        p = mtf.paragraphs[0]
        p.text = url
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(200, 230, 255)
        p.alignment = PP_ALIGN.LEFT
        meta_top += line_gap

    # --- optional watermark / brand bg on right side ------------------------
    brand_bg = assets.get("brand_bg")
    if brand_bg and os.path.exists(brand_bg):
        try:
            slide.shapes.add_picture(
                brand_bg,
                Inches(content_left),
                Inches(sh_in - 1.2),
                height=Inches(1.0),
            )
        except Exception:
            # ignore image errors; we don't want to break slide creation
            pass



# ----------------------------
# Build
# ----------------------------

def build(
    inputs: List[str],
    output_path: str,
    month: Optional[str],
    assets: dict,
    template: Optional[str],
    rail_width: float,
    conclusion_links: Optional[List[tuple[str, str]]] = None,
):
    prs = Presentation(template) if (template and os.path.exists(template)) else Presentation()

    # Core section slides
    add_cover_slide(
        prs,
        assets,
        assets.get("cover_title"),
        assets.get("cover_dates"),
        assets.get("logo"),
        assets.get("logo2"),
    )
    add_agenda_slide(prs, assets)

    # 1) Aggregate
    all_items: list[Item] = []
    for path in inputs:
        dicts = (parse_message_center_html(path, month)
                if "messagecenter" in path.lower() or "briefing" in path.lower()
                else parse_roadmap_html(path, month))
        all_items.extend(_to_item(d) for d in dicts)

    # 2) Dedup/sort -> this becomes the canonical 'items'
    def _dedup_keep_order(seq: list[Item]) -> list[Item]:
        seen: set[str] = set()
        out: list[Item] = []
        for it in seq:
            key = (it.roadmap_id or it.title or it.url or "").lower().strip()
            if not key or key in seen:
                continue
            seen.add(key)
            out.append(it)
        return out

    items = _dedup_keep_order(all_items)
    items.sort(key=lambda i: (i.products or "", i.title or ""))

    # 3) Everything below should reference `items`
    # Optional: per-product separators before items
    
    if items:
        # Group by product (stable order by first appearance)
        order: List[str] = []
        by_product: Dict[str, List[Item]] = {}
        for it in items:
            prod = (it.products[0] if it.products else "General").strip() or "General"
            if prod not in by_product:
                by_product[prod] = []
                order.append(prod)
            by_product[prod].append(it)

        # Build slides
        idx = 1
        for prod in order:
            add_separator_slide(prs, assets, title=f"{prod} updates", subtitle=month or "")
            for it in by_product[prod]:
                add_item_slide(
                    prs,
                    it,
                    month_str=month or "",
                    assets=assets,
                    rail_left_in=rail_width,
                    rail_width_in=rail_width,
                )
                idx += 1
    else:
        # Fallback if no items
        add_separator_slide(prs, assets, title="No updates found", subtitle=month or "")

    # Conclusion + Thank you
    if not conclusion_links:
        conclusion_links = [
            ("Microsoft Security", "https://www.microsoft.com/en-us/security"),
            ("Azure Updates", "https://azure.microsoft.com/en-us/updates/"),
            (
                "Dynamics 365 & Power Platform",
                "https://www.microsoft.com/en-us/dynamics-365/?culture=en-us&country=us",
            ),
            (
                "Technical Documentation",
                "https://learn.microsoft.com/en-us/docs/?culture=en-us&country=us",
            ),
        ]
    add_conclusion_slide(prs, assets, conclusion_links)
    add_thankyou_slide(prs, assets)

    actual = _safe_save(prs, output_path)
    print(f"[ok] Deck saved to: {actual}")



# ----------------------------
# CLI
# ----------------------------
def main():
    p = argparse.ArgumentParser()
    p.add_argument(
        "-i",
        "--inputs",
        nargs="+",
        required=True,
        help="One or more HTML inputs (Roadmap/Message Center)",
    )
    p.add_argument("-o", "--output", required=True, help="Output .pptx")
    p.add_argument("--month", default="", help="Month label like 'September 2025'")
    # assets
    p.add_argument("--cover", dest="cover", default="", help="Cover background image")
    p.add_argument("--agenda-bg", dest="agenda_bg", default="", help="Agenda background image")
    p.add_argument("--separator", dest="separator", default="", help="Separator background image")
    p.add_argument(
        "--conclusion-bg", dest="conclusion_bg", default="", help="Conclusion background image"
    )
    p.add_argument("--thankyou", dest="thankyou", default="", help="Thank-you background image")
    p.add_argument(
        "--brand-bg", dest="brand_bg", default="", help="Generic brand background for item slides"
    )
    p.add_argument("--cover-title", dest="cover_title", default="", help="Cover title text")
    p.add_argument("--cover-dates", dest="cover_dates", default="", help="Cover dates text")
    p.add_argument(
        "--separator-title",
        dest="separator_title",
        default="",
        help="(unused; separators auto from product)",
    )
    p.add_argument("--logo", dest="logo", default="", help="Logo 1 (e.g., Parex)")
    p.add_argument("--logo2", dest="logo2", default="", help="Logo 2 (e.g., customer)")
    p.add_argument("--rocket", dest="rocket", default="", help="Path to rocket icon")
    p.add_argument("--magnifier", dest="magnifier", default="", help="Path to magnifier icon")
    p.add_argument("--admin", dest="admin", default="", help="Admin target-audience icon")
    p.add_argument("--user", dest="user", default="", help="User target-audience icon")
    p.add_argument("--check", dest="check", default="", help="Green check icon")
    p.add_argument("--template", dest="template", default="", help="Optional template .pptx")
    p.add_argument(
        "--rail-width",
        dest="rail_width",
        default=str(DEFAULT_RAIL_WIDTH_IN),
        help="Right rail width in inches (default 3.5)",
    )

    args = p.parse_args()

    assets = {
        "cover": path_if_exists(args.cover),
        "agenda": path_if_exists(args.agenda_bg),
        "separator": path_if_exists(args.separator),
        "conclusion": path_if_exists(args.conclusion_bg),
        "thankyou": path_if_exists(args.thankyou),
        "brand_bg": path_if_exists(args.brand_bg),
        "cover_title": args.cover_title or "M365 Technical Update Briefing",
        "cover_dates": args.cover_dates or args.month or "",
        "logo": path_if_exists(args.logo),
        "logo2": path_if_exists(args.logo2),
        "rocket": path_if_exists(args.rocket)
        or path_if_exists(os.path.join("assets", "rocket.png")),
        "magnifier": path_if_exists(args.magnifier)
        or path_if_exists(os.path.join("assets", "magnifier.png")),
        "admin": path_if_exists(args.admin) or path_if_exists(os.path.join("assets", "admin.png")),
        "user": path_if_exists(args.user) or path_if_exists(os.path.join("assets", "user.png")),
        "check": path_if_exists(args.check) or path_if_exists(os.path.join("assets", "check.png")),
    }

    rail_w = float(args.rail_width) if args.rail_width else DEFAULT_RAIL_WIDTH_IN

    build(
        inputs=args.inputs,
        output_path=args.output,
        month=args.month or "",
        assets=assets,
        template=args.template or None,
        rail_width=rail_w,
        conclusion_links=None,
    )


if __name__ == "__main__":
    main()
