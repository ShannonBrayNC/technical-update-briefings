from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def inches_to_emu(x):
    return int(x * 914400)

def add_full_bg(slide, img_path):
    # Just a placeholder, replace with actual background filling if needed.
    pass

def add_title_box(slide, text, *, left_in, top_in, width_in, height_in, font_size_pt=60, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER):
    from pptx.util import Inches
    textbox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size_pt)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)

def add_text_box(slide, text, *, left_in, top_in, width_in, height_in, font_size_pt=16, color="#FFFFFF", bold=False, align=PP_ALIGN.LEFT):
    from pptx.util import Inches
    textbox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size_pt)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)

class Layout:
    def __init__(self, assets):
        self.assets = assets

    def add_cover_slide(self, prs, assets, content):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_bg(slide, assets.get("cover"))
        add_title_box(
            slide,
            content.get('title', 'Title'),
            left_in=0.5, top_in=2, width_in=9, height_in=2,
            font_size_pt=52, color="#FFFFFF", bold=True
        )
        add_text_box(
            slide,
            content.get('dates', ''),
            left_in=0.5, top_in=4, width_in=9, height_in=1,
            font_size_pt=30, color="#FFFFFF"
        )

    def add_item_slide(self, prs, item, month_str, assets, rail_left_in=0):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Fill background if needed
        add_full_bg(slide, assets.get("cover"))  # or custom background

        # Title
        add_title_box(slide, item.title or 'Title', left_in=0.5, top_in=0.5, width_in=9, height_in=1.8, font_size_pt=36)
        # Summary
        add_text_box(slide, item.summary or '', left_in=0.5, top_in=2.5, width_in=9, height_in=3, font_size_pt=16)

        # Meta info (e.g., feature_id, status)
        add_text_box(slide, f"ID: {item.roadmap_id}", left_in=0.5, top_in=6, width_in=3, height_in=0.5, font_size_pt=12)
        add_text_box(slide, f"Status: {item.status}", left_in=3.5, top_in=6, width_in=3, height_in=0.5, font_size_pt=12)

    def add_separator_slide(self, prs, assets, title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_bg(slide, assets.get("separator"))
        add_title_box(slide, title, left_in=0.5, top_in=1.0, width_in=9, height_in=1.5, font_size_pt=48)

    def add_conclusion_slide(self, prs, assets, links):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_bg(slide, assets.get("conclusion"))
        add_title_box(slide, "Final Thoughts", left_in=0.5, top_in=1.0, width_in=9, height_in=1.5, font_size_pt=48)
        top = 3
        for text, url in links:
            add_text_box(slide, f"{text}: {url}", left_in=0.5, top_in=top, width_in=9, height_in=0.5, font_size_pt=16)
            top += 0.6

    def add_thankyou_slide(self, prs, assets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_bg(slide, assets.get("thankyou"))
