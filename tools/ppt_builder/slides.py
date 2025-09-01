# slides.py — unified: canonical API + "pro" helpers
import os
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# -------------------- basics --------------------

def inches_to_emu(x: float) -> int:
    return int(x * 914400)

def _hex(color: str) -> str:
    s = str(color or "FFFFFF").strip()
    if s.startswith("#"):
        s = s[1:]
    if len(s) == 3 and all(c in "0123456789aAbBcCdDeEfF" for c in s):
        s = "".join(c*2 for c in s)
    if len(s) != 6:
        s = "FFFFFF"
    return s.upper()

# -------------------- "pro" helpers (keep names) --------------------

def add_full_bg(slide, image_path: str):
    """Full-bleed bg image; falls back to a subtle dark fill if missing."""
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, width=inches_to_emu(10), height=inches_to_emu(7.5))
        return
    # fallback: dark rectangle
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inches_to_emu(10), inches_to_emu(7.5))
    fill = shp.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(22, 28, 36)  # #161C24
    shp.line.fill.background()

def add_title(slide, text, *, top=0.9, size=52, align=PP_ALIGN.LEFT, color="FFFFFF", name="Segoe UI Semibold"):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(1.2))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text or ""
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = name
    r.font.color.rgb = RGBColor.from_string(_hex(color))

def add_text(slide, text, *, left=0.6, top=2.0, width=8.8, height=3.4, size=20, color="FFFFFF", name="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text or ""
    r.font.size = Pt(size)
    r.font.bold = False
    r.font.name = name
    r.font.color.rgb = RGBColor.from_string(_hex(color))

def add_chip(slide, text, *, left, top, fill="2E7D32", text_color="FFFFFF"):
    """Rounded chip for statuses like 'Preview', 'GA' etc."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches_to_emu(left), inches_to_emu(top), inches_to_emu(1.6), inches_to_emu(0.45))
    try:
        shp.adjustments[0] = 0.5  # more rounded, when available
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(_hex(fill))
    shp.line.fill.background()
    tb = shp.text_frame
    tb.clear(); tb.word_wrap = True
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text or ""
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(_hex(text_color))

def add_picture_safe(slide, image_path: str, left_in: float, top_in: float, width_in=None, height_in=None):
    if not image_path or not os.path.exists(image_path):
        return
    if width_in is None and height_in is None:
        slide.shapes.add_picture(image_path, inches_to_emu(left_in), inches_to_emu(top_in))
    elif width_in is not None and height_in is None:
        slide.shapes.add_picture(image_path, inches_to_emu(left_in), inches_to_emu(top_in), width=inches_to_emu(width_in))
    elif width_in is None and height_in is not None:
        slide.shapes.add_picture(image_path, inches_to_emu(left_in), inches_to_emu(top_in), height=inches_to_emu(height_in))
    else:
        slide.shapes.add_picture(image_path, inches_to_emu(left_in), inches_to_emu(top_in), width=inches_to_emu(width_in), height=inches_to_emu(height_in))

# -------------------- canonical API (wrappers over pro helpers) --------------------

def add_full_slide_picture(slide, prs, image_path: str):
    """Canonical wrapper used by run_build; delegates to 'add_full_bg'."""
    add_full_bg(slide, image_path)

def add_title_box(slide, text: str, *, left_in: float, top_in: float, width_in: float, height_in: float,
                  font_size_pt: int = 60, bold: bool = True, color: str = "FFFFFF", align=PP_ALIGN.LEFT):
    """Canonical wrapper; positions like old API but uses 'add_title' styling."""
    # honor left/width by creating frame at those coords
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text or ""
    r.font.size = Pt(font_size_pt)
    r.font.bold = bool(bold)
    r.font.name = "Segoe UI Semibold"
    r.font.color.rgb = RGBColor.from_string(_hex(color))

def add_text_box(slide, text: str, *, left_in: float, top_in: float, width_in: float, height_in: float,
                 font_size_pt: int = 18, bold: bool = False, color: str = "FFFFFF", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text or ""
    r.font.size = Pt(font_size_pt)
    r.font.bold = bool(bold)
    r.font.name = "Segoe UI"
    r.font.color.rgb = RGBColor.from_string(_hex(color))

def add_cover_slide(prs, assets: dict, cover_title: str, cover_dates: str, logo1_path: str, logo2_path: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("cover"))
    add_title(slide, cover_title or "M365 Technical Update Briefing", top=2.1, size=54, align=PP_ALIGN.CENTER)
    if cover_dates:
        add_text(slide, cover_dates, left=0.6, top=3.8, width=8.8, height=0.8, size=30, color="E6E8EF")
    if logo1_path:
        add_picture_safe(slide, logo1_path, left_in=0.4, top_in=6.6, height_in=0.6)
    if logo2_path:
        add_picture_safe(slide, logo2_path, left_in=8.0, top_in=6.6, height_in=0.6)
    if cover_title:
        add_picture_safe(slide, cover_title, left_in=0.4, top_in=6.6, height_in=0.6)
           

def add_agenda_slide(prs, assets: dict, agenda_lines=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("agenda"))
    add_title(slide, "Agenda", top=0.9, size=52, align=PP_ALIGN.LEFT)
    if not agenda_lines:
        agenda_lines = ["Overview", "Key updates by product", "Timeline & rollout status", "Q&A"]
    top = 2.1
    for line in agenda_lines:
        add_text(slide, f"• {line}", left=0.9, top=top, width=8.2, height=0.5, size=26, color="FFFFFF")
        top += 0.55

def add_separator_slide(prs, assets: dict, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("separator"))
    add_title(slide, title, top=3.0, size=56, align=PP_ALIGN.CENTER)

def add_conclusion_slide(prs, assets: dict, links):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, assets.get("conclusion"))
    add_title(slide, "Final Thoughts", top=0.9, size=52, align=PP_ALIGN.CENTER)
    top = 2.3
    for text, url in (links or []):
        add_text(slide, f"{text}: {url}", left=0.6, top=top, width=8.8, height=0.6, size=22, color="E6E8EF")
        top += 0.6
