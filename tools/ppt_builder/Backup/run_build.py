# run_build.py — full builder (with rail card, product palette, icons, auto title-height)

from __future__ import annotations

import os
import json
import re
import difflib
import importlib.util
import textwrap
from typing import Dict, List, Optional, Tuple, cast

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

import slides as S  # relies on add_* helpers

# ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
PARSERS_PATH = os.path.join(HERE, "parsers.py")

def _log(msg: str) -> None:
    print(f"[run_build] {msg}", flush=True)

# Dynamic import of parsers
spec = importlib.util.spec_from_file_location("parsers_local", PARSERS_PATH)
if not spec or not spec.loader:
    raise ImportError(f"Cannot load parsers.py at {PARSERS_PATH}")
P = importlib.util.module_from_spec(spec)  # type: ignore[arg-type]
spec.loader.exec_module(P)  # type: ignore[union-attr]

parse_message_center_html = P.parse_message_center_html
parse_roadmap_html = P.parse_roadmap_html

# ------------------------------------------------------------------
# Parsing & merge
def _safe_parse(path: str, month: str) -> List[dict]:
    low = (path or "").lower()
    is_mc = ("messagecenter" in low) or ("message_center" in low) or ("briefing" in low)
    fn = parse_message_center_html if is_mc else parse_roadmap_html
    try:
        items = fn(path, month) or []
        for it in items:
            it.setdefault("source", "mc" if is_mc else "rm")
            it.setdefault("title",""); it.setdefault("products",[])
            it.setdefault("clouds",[]); it.setdefault("platforms",[])
            it.setdefault("summary",""); it.setdefault("description","")
        _log(f"Parsed {len(items)} with month='{month}' from {path}")
        if len(items) == 0 and month:
            _log("No items after month filter — retrying with no month filter...")
            items = fn(path, "") or []
            for it in items:
                it.setdefault("source", "mc" if is_mc else "rm")
            _log(f"Parsed {len(items)} with month='' from {path}")
        return items
    except Exception as e:
        _log(f"ERROR parsing {path}: {e}")
        return []

def _norm_title(t: str) -> str:
    t = re.sub(r"\s+", " ", (t or "").strip()).lower()
    t = re.sub(r"[^a-z0-9 ]+", "", t)
    return t

def _sim(a: str, b: str) -> float:
    return difflib.SequenceMatcher(None, _norm_title(a), _norm_title(b)).ratio()

def _merge_record(base: dict, other: dict) -> dict:
    def pick_text(a: Optional[str], b: Optional[str]) -> str:
        a = a or ""; b = b or ""; return b if len(b) > len(a) else a

    src_base, src_other = base.get("source",""), other.get("source","")
    mc_first = (src_other == "mc") and (src_base != "mc")

    if mc_first:
        base["title"] = other.get("title") or base.get("title") or ""
        base["summary"] = pick_text(base.get("summary"), other.get("summary"))
        base["description"] = pick_text(base.get("description"), other.get("description"))
        base["status"] = other.get("status") or base.get("status") or ""
        base["ga"] = other.get("ga") or base.get("ga") or ""
        if other.get("url"): base["url"] = other["url"]
    else:
        base["summary"] = pick_text(base.get("summary"), other.get("summary"))
        base["description"] = pick_text(base.get("description"), other.get("description"))
        if not base.get("status"): base["status"] = other.get("status","")
        if not base.get("ga"):     base["ga"]     = other.get("ga","")
        if not base.get("url"):    base["url"]    = other.get("url","")

    for k in ("products","clouds","platforms","audience"):
        a = base.get(k) or []
        b = other.get(k) or []
        base[k] = sorted({*(x for x in a if x), *(y for y in b if y)}, key=str.lower)

    if not base.get("roadmap_id"):
        base["roadmap_id"] = other.get("roadmap_id","")

    return base

def _merge_items(items: List[dict]) -> List[dict]:
    by_id: Dict[str, dict] = {}
    no_id: List[dict] = []
    for it in items:
        rid = (it.get("roadmap_id") or "").strip()
        (by_id if rid else no_id).setdefault(rid, dict(it)) if rid and rid not in by_id else (
            by_id.update({rid: _merge_record(by_id[rid], it)}) if rid else no_id.append(it)
        )

    merged: List[dict] = list(by_id.values())

    def overlap(a: Optional[List[str]], b: Optional[List[str]]) -> bool:
        return bool(set(a or []).intersection(set(b or [])))

    for it in no_id:
        best_idx, best_score = -1, 0.0
        for idx, rec in enumerate(merged):
            if not overlap(it.get("products"), rec.get("products")): continue
            s = _sim(it.get("title",""), rec.get("title",""))
            if s > best_score: best_idx, best_score = idx, s
        if best_idx >= 0 and best_score >= 0.82:
            merged[best_idx] = _merge_record(merged[best_idx], it)
        else:
            merged.append(dict(it))

    out: List[dict] = []
    for it in merged:
        if not out: out.append(it); continue
        best_idx, best_score = -1, 0.0
        for idx, rec in enumerate(out):
            if overlap(it.get("products"), rec.get("products")):
                s = _sim(it.get("title",""), rec.get("title",""))
                if s > best_score: best_idx, best_score = idx, s
        if best_idx >= 0 and best_score >= 0.85:
            out[best_idx] = _merge_record(out[best_idx], it)
        else:
            out.append(it)
    return out

# ------------------------------------------------------------------
# Rendering helpers
def _titlecase(s: str) -> str:
    if not s: return ""
    small = {"and","or","for","nor","a","an","the","as","at","by","in","of","on","per","to","vs","via"}
    words = s.strip().split()
    out: List[str] = []
    for i, w in enumerate(words):
        if w.isupper() or any(ch.isdigit() for ch in w): out.append(w); continue
        lw = w.lower()
        out.append(lw if i not in (0, len(words)-1) and lw in small else (w[:1].upper() + w[1:].lower()))
    return " ".join(out)

def _add_rail(slide, rail_width_in: float = 3.5, hex_color: str = "0F172A") -> None:
    EMU = 914400
    rail_w = int(rail_width_in * EMU)
    full_w = int(10 * EMU)
    full_h = int(7.5 * EMU)
    left = full_w - rail_w
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, 0, rail_w, full_h)
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    shp.line.fill.background()

def _rail_color_for(products: Optional[List[str]], assets: Dict) -> str:
    # prefer config palette
    palette = { (k or "").lower(): (v or "") for k,v in (assets.get("product_palette") or {}).items() }
    p = (products or [""])[0].lower()
    if p in palette and palette[p]:
        return palette[p]

    # fallback defaults
    default = {
        "teams":      "5B21B6",
        "sharepoint": "065F46",
        "onedrive":   "1D4ED8",
        "exchange":   "0F172A",
        "security":   "14532D",
        "outlook":    "1F2937",
        "defender":   "14532D",
        "purview":    "065F46",
        "entra":      "1D4ED8",
    }
    return default.get(p, (assets.get("theme",{}).get("rail_default_hex") or "0F172A"))

def _add_footer(slide, month_label: str, page: Optional[int], total: Optional[int]) -> None:
    left = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(4.0), Inches(0.35))
    tf = left.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = month_label or ""
    r.font.size = Pt(10); r.font.color.rgb = RGBColor.from_string("94A3B8")
    if page is not None and total is not None:
        right = slide.shapes.add_textbox(Inches(9.2), Inches(7.05), Inches(0.9), Inches(0.35))
        tf2 = right.text_frame; tf2.clear()
        p2 = tf2.paragraphs[0]; r2 = p2.add_run()
        r2.text = f"{page}/{total}"
        r2.font.size = Pt(10); r2.font.color.rgb = RGBColor.from_string("94A3B8")

def _status_icon_for(status: Optional[str], assets: Dict) -> Optional[str]:
    s = (status or "").lower()
    if any(k in s for k in ("ga", "rolling out", "rollout")):
        return cast(Optional[str], assets.get("icon_rocket") or "")
    if any(k in s for k in ("preview", "in development", "dev")):
        return cast(Optional[str], assets.get("icon_preview") or "")
    return None

def _rail_card(slide, rail_left_in: float, rail_w_in: float, month: str, rid: str) -> None:
    if rail_w_in <= 0: return
    # card
    card_left = rail_left_in + 0.25
    card_top  = 1.6
    card_w, card_h = max(rail_w_in - 0.5, 1.4), 3.0
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(card_left), Inches(card_top),
                                 Inches(card_w), Inches(card_h))
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    shp.line.fill.background()

    # month (headline)
    t = slide.shapes.add_textbox(Inches(card_left+0.2), Inches(card_top+0.25),
                                 Inches(card_w-0.4), Inches(0.7)).text_frame
    p = t.paragraphs[0]; r = p.add_run()
    r.text = month or ""; r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("6B21A8")

    # Roadmap ID
    t2 = slide.shapes.add_textbox(Inches(card_left+0.2), Inches(card_top+1.2),
                                  Inches(card_w-0.4), Inches(0.45)).text_frame
    p2 = t2.paragraphs[0]; r2 = p2.add_run()
    r2.text = f"Roadmap ID: {rid or '—'}"; r2.font.size = Pt(14)
    r2.font.color.rgb = RGBColor.from_string("111827")

# ------------------------------------------------------------------
def _build_item_slide(
    prs,
    item: Dict,
    month: str,
    assets: Dict,
    rail_width: Optional[float] = None,
    page: Optional[int] = None,
    total: Optional[int] = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # background
    bg_path: str = cast(str, assets.get("brand_bg") or assets.get("cover") or "")
    if bg_path:
        S.add_full_slide_picture(slide, prs, bg_path)

    rail_w = float(rail_width or 0.0)
    if rail_w > 0:
        rail_hex = _rail_color_for(item.get("products"), assets)
        _add_rail(slide, rail_width_in=rail_w, hex_color=rail_hex)

    rail_left = 10.0 - rail_w

    # Title (auto height) and Body start
    raw_title = item.get("title", "") or item.get("headline","")
    title = _titlecase(raw_title)

    w_body = 10.0 - rail_w - 0.9
    TITLE_TOP, TITLE_FONT_PT = 0.6, 34
    LINE_HEIGHT_IN, MIN_TITLE_H, MAX_TITLE_H = 0.55, 1.10, 1.60

    wrap_cols = max(20, int(w_body * 4.0))
    est_lines = max(1, len(textwrap.wrap(title, width=wrap_cols)))
    title_h   = min(MAX_TITLE_H, max(MIN_TITLE_H, est_lines * LINE_HEIGHT_IN))

    S.add_title_box(
        slide, title,
        left_in=0.6, top_in=TITLE_TOP, width_in=w_body, height_in=title_h,
        font_size_pt=TITLE_FONT_PT, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER
    )

    body_top = TITLE_TOP + title_h + 0.30

    # Status chip + icon (centered in rail)
    status = (item.get("status") or "").lower()
    chip_w = 1.8
    chip_left = rail_left + max(0.1, (rail_w - chip_w) / 2.0) if rail_w > 0 else (0.6 + w_body - chip_w)
    color_map = {
        "ga": ("16A34A","FFFFFF"),
        "general": ("16A34A","FFFFFF"),
        "rolling out": ("3B82F6","FFFFFF"),
        "preview": ("F59E0B","111827"),
        "in development": ("6B7280","FFFFFF"),
    }
    for key, (fill, fg) in color_map.items():
        if key in status:
            S.add_chip(slide, status.title(), left=chip_left, top=0.60, fill=fill, text_color=fg)
            break

    icon_path = _status_icon_for(status, assets)
    if icon_path and rail_w > 0:
        icon_w = 1.1; icon_h = 1.1
        icon_left = rail_left + max(0.1, (rail_w - icon_w) / 2.0)
        slide.shapes.add_picture(icon_path, Inches(icon_left), Inches(0.35), width=Inches(icon_w), height=Inches(icon_h))

    # Right-rail card (Month + Roadmap ID)
    if rail_w > 0:
        _rail_card(slide, rail_left, rail_w, month, item.get("roadmap_id",""))

    # Body text
    body = item.get("summary") or item.get("description") or item.get("body") or ""
    if not body:
        bullets: List[str] = []
        if item.get("status"):    bullets.append(f"• Status: {item['status']}")
        if item.get("ga"):        bullets.append(f"• GA: {item['ga']}")
        if item.get("platforms"): bullets.append("• Platforms: " + ", ".join(item["platforms"]))
        if item.get("clouds"):    bullets.append("• Clouds: " + ", ".join(item["clouds"]))
        body = "\n".join(bullets)

    S.add_text_box(
        slide, body,
        left_in=0.6, top_in=body_top, width_in=w_body, height_in=3.9,
        font_size_pt=20, bold=False, color="FFFFFF"
    )

    # Meta footer (no raw URL here)
    meta_lines: List[str] = []
    for k in ("roadmap_id","status","ga"):
        v = item.get(k)
        if v: meta_lines.append(f"{k.upper()}: {v}")
    if item.get("products"): meta_lines.append("Product: " + ", ".join(item["products"]))
    if item.get("clouds"):   meta_lines.append("Clouds: " + ", ".join(item["clouds"]))
    if month:                meta_lines.append(month)

    S.add_text_box(
        slide, "  •  ".join(meta_lines),
        left_in=0.6, top_in=6.2, width_in=w_body, height_in=0.7,
        font_size_pt=14, bold=False, color="E6E8EF"
    )

    # Clickable link
    url = item.get("url")
    if url:
        link = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(2.2), Inches(0.35))
        tf = link.text_frame; tf.clear(); tf.word_wrap = False
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = "Open item ↗"; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = RGBColor(59, 130, 246); r.hyperlink.address = url

    _add_footer(slide, month, page, total)

# ------------------------------------------------------------------
def build(
    inputs: List[str],
    output_path: str,
    month: str,
    assets: Dict,
    template: str | None = None,
    rail_width=None,
    conclusion_links=None,
    debug_dump: Optional[str] = None,
) -> None:
    _log("Starting deck build...")

    all_items: List[dict] = []
    for fp in inputs:
        all_items.extend(_safe_parse(fp, month))

    _log(f"Raw items before merge: {len(all_items)}")
    merged_items = _merge_items(all_items)
    _log(f"After MC-first merge + fuzzy: {len(merged_items)} unique items")

    if debug_dump:
        try:
            os.makedirs(os.path.dirname(os.path.abspath(debug_dump)) or ".", exist_ok=True)
            with open(debug_dump, "w", encoding="utf-8") as f:
                json.dump(merged_items, f, indent=2, ensure_ascii=False)
            _log(f"Debug dump wrote {len(merged_items)} items to {os.path.abspath(debug_dump)}")
        except Exception as e:
            _log(f"Failed to write debug dump: {e}")

    def first_prod(it): return (it.get("products") or ["General"])[0].lower()
    merged_items.sort(key=lambda i: (first_prod(i), (i.get("title") or '').lower()))

    prs = Presentation(template) if (template and os.path.exists(template)) else Presentation()

    cover_title: str = cast(str, assets.get("cover_title") or "M365 Technical Update Briefing")
    cover_dates: str = cast(str, assets.get("cover_dates") or (month or ""))
    logo1: str = cast(str, assets.get("logo") or "")
    logo2: str = cast(str, assets.get("logo2") or "")
    S.add_cover_slide(prs, assets, cover_title, cover_dates, logo1, logo2)
    S.add_agenda_slide(prs, assets)

    grouped: Dict[str, List[dict]] = {}
    order: List[str] = []
    for it in merged_items:
        p = (it.get("products") or ["General"])[0]
        if p not in grouped:
            grouped[p] = []; order.append(p)
        grouped[p].append(it)

    item_total = sum(len(v) for v in grouped.values())
    page_counter = 0

    if not order:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        S.add_title_box(slide, "No updates parsed",
                        left_in=0.6, top_in=2.0, width_in=8.8, height_in=1.2,
                        font_size_pt=40, color="FFFFFF")
    else:
        for prod in order:
            S.add_separator_slide(prs, assets, f"{prod} updates")
            for it in grouped[prod]:
                page_counter += 1
                _build_item_slide(
                    prs, it, month, assets,
                    rail_width=float(rail_width) if rail_width else None,
                    page=page_counter, total=item_total
                )

    links = conclusion_links or [("Security","https://www.microsoft.com/security"),
                                 ("Azure","https://azure.microsoft.com/"),
                                 ("Docs","https://learn.microsoft.com/")]
    S.add_conclusion_slide(prs, assets, links)
    thank = cast(str, assets.get("thankyou") or "")
    if thank:
        S.add_full_slide_picture(prs.slides.add_slide(prs.slide_layouts[6]), prs, thank)

    out_abs = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(out_abs) or ".", exist_ok=True)
    prs.save(out_abs)
    _log(f"Deck saved: {out_abs}")

# ------------------------------------------------------------------
if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("-i","--inputs", nargs="+", required=True)
    ap.add_argument("-o","--output", required=True)
    ap.add_argument("--month", default="")
    ap.add_argument("--template", default="")
    ap.add_argument("--cover", default="")
    ap.add_argument("--agenda", default="")
    ap.add_argument("--separator", default="")
    ap.add_argument("--conclusion", default="")
    ap.add_argument("--thankyou", default="")
    ap.add_argument("--brand_bg", dest="brand_bg", default="")
    ap.add_argument("--cover_title", default="M365 Technical Update Briefing")
    ap.add_argument("--cover_dates", default="")
    ap.add_argument("--logo", default="")
    ap.add_argument("--logo2", default="")
    ap.add_argument("--rail_width", type=float, default=None)
    ap.add_argument("--icon_rocket", dest="icon_rocket", default="")
    ap.add_argument("--icon_preview", dest="icon_preview", default="")
    ap.add_argument("--debug_dump", default="")
    args = ap.parse_args()

    assets = {
        "cover": args.cover,
        "agenda": args.agenda,
        "separator": args.separator,
        "conclusion": args.conclusion,
        "thankyou": args.thankyou,
        "brand_bg": args.brand_bg,
        "cover_title": args.cover_title,
        "cover_dates": args.cover_dates,
        "logo": args.logo,
        "logo2": args.logo2,
        "icon_rocket": args.icon_rocket,
        "icon_preview": args.icon_preview,
    }

    build(
        args.inputs, args.output, args.month, assets,
        template=args.template, rail_width=args.rail_width,
        conclusion_links=None, debug_dump=(args.debug_dump or None)
    )
