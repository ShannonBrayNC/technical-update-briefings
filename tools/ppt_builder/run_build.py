#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
run_build.py
------------
Builds a styled PowerPoint deck from HTML inputs (Roadmap + Message Center).
- MC-first merge with fuzzy dedupe
- Right rail with product-colored background, big status chip, status icon
- Rail card (Month, Roadmap ID, compact details)
- "Target Audience" header + two rows (End Users, Admins)
- Left side: centered title with auto height clamp; summary paragraph (no bullets)
- Background comes from brand_bg (else cover)
- Footer shows month + page/total

CLI (underscored flags) for direct invocation.
Use generate_deck.py (hyphenated flags) if you need style.yaml support.
"""

from __future__ import annotations

import argparse
import importlib
import json
import re
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Tuple, Union, Iterable

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ---- slides helpers (local) -------------------------------------------------
# Expect a local slides.py providing add_full_slide_picture/add_title_box/add_text_box.
# We import safely and fail clearly if missing.
import importlib.util as _impspec
_slides_mod = None
for cand in ("slides.py", "deck_slides.py"):
    p = Path(__file__).with_name(cand)
    if p.exists():
        spec = _impspec.spec_from_file_location("slides", str(p))
        if spec and spec.loader:
            _slides_mod = _impspec.module_from_spec(spec)
            spec.loader.exec_module(_slides_mod)
            break
if _slides_mod is None:
    raise ImportError("Local slides.py/deck_slides.py not found or failed to import.")
S = _slides_mod


def _tenant_from_clouds(item: Dict) -> str:
    clouds = item.get("clouds")
    if isinstance(clouds, list):
        s = ", ".join(str(x) for x in clouds if x)
    else:
        s = str(clouds or "")
    low = s.lower()
    if "dod" in low: return "DoD"
    if "gcc high" in low: return "GCC High"
    if "gcc" in low or "government" in low: return "GCC"
    if "worldwide" in low or "standard multi-tenant" in low or "multi-tenant" in low:
        return "Worldwide (Standard Multi-Tenant)"
    return s.strip()


from datetime import datetime

def _fmt_month(s: str) -> str:
    s = (s or "").strip()
    if not s: return ""
    # try YYYY-MM-DD or YYYY-MM
    for fmt in ("%Y-%m-%d", "%Y-%m"):
        try:
            d = datetime.strptime(s[:len(fmt)], fmt)
            return d.strftime("%B %Y")
        except Exception:
            pass
    # already looks like "September 2025"
    return s

def _phase_text(item: Dict) -> str:
    p = item.get("phases")
    if isinstance(p, (list, tuple, set)): p = ", ".join([str(x).strip() for x in p if x])
    return (p or "").strip()

def _rail_card(slide, rail_left_in: float, rail_w_in: float, month_label: str, item: Dict) -> tuple[float, float]:
    card_top_in = 1.10
    card_h_in   = 2.05
    pad_in      = 0.30

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(rail_left_in + 0.25),
                                 Inches(card_top_in),
                                 Inches(rail_w_in - 0.50),
                                 Inches(card_h_in))
    shp.adjustments[0] = 0.18
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string("F8FAFC")
    shp.line.color.rgb = RGBColor.from_string("0B1220"); shp.line.width = Pt(1)

    tf = shp.text_frame; tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(pad_in)

    rollout = _fmt_month(item.get("rollout_start")) or _fmt_month(item.get("ga_start")) \
              or _fmt_month(item.get("ga_end")) or (month_label or "")

    # Month / Year (rollout)
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = rollout; r.font.bold = True; r.font.size = Pt(20)
    r.font.color.rgb = RGBColor.from_string("6B21A8")

    # Status
    status = (item.get("status") or "").strip()
    p = tf.add_paragraph(); p.space_before = Pt(2)
    r1 = p.add_run(); r1.text = "Status: "; r1.font.bold = True; r1.font.size = Pt(12)
    r1.font.color.rgb = RGBColor.from_string("111827")
    r2 = p.add_run(); r2.text = status.title() if status else "—"; r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor.from_string("374151")

    # Phase(s)
    phases = _phase_text(item)
    p = tf.add_paragraph()
    r1 = p.add_run(); r1.text = "Phase(s): "; r1.font.bold = True; r1.font.size = Pt(12)
    r1.font.color.rgb = RGBColor.from_string("111827")
    r2 = p.add_run(); r2.text = phases or "—"; r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor.from_string("374151")

    # Roadmap ID
    rid = (item.get("roadmap_id") or "").strip()
    p = tf.add_paragraph()
    r1 = p.add_run(); r1.text = "Roadmap ID: "; r1.font.bold = True; r1.font.size = Pt(12)
    r1.font.color.rgb = RGBColor.from_string("111827")
    r2 = p.add_run(); r2.text = rid if rid else "—"; r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor.from_string("374151")

    return card_top_in, card_h_in


# ---- logging ----------------------------------------------------------------
def _log(msg: str) -> None:
    print(f"[run_build] {msg}", flush=True)

# ---- simple HTML parsers (prefer external 'parsers' if available) -----------+
def _try_external_parsers() -> Optional[object]:
    try:
        spec = _impspec.find_spec("parsers")
        if spec is None:
            return None
        mod = importlib.import_module("parsers")  # type: ignore
        return mod
    except Exception:
        return None

def _parse_one(path: str, month: str) -> List[Dict]:
    """Parse a single HTML file into a list of item dicts.
    Tries external 'parsers' first; else uses a forgiving built-in fallback.
    """
    mod = _try_external_parsers()
    p = Path(path)
    src = "message_center" if "message_center" in str(p).lower() else "roadmap"
    items: List[Dict] = []
    if mod:
        try:
            if src == "message_center" and hasattr(mod, "parse_message_center_html"):
                items = mod.parse_message_center_html(str(p), month)  # type: ignore
            elif src == "roadmap" and hasattr(mod, "parse_roadmap_html"):
                items = mod.parse_roadmap_html(str(p), month)  # type: ignore
        except Exception as e:
            _log(f"ERROR external parser on {path}: {e}")

    if items:
        _log(f"Parsed {len(items)} with month='{month}' from {path}")
        for it in items:
            it.setdefault("_source", src)
        return items

    # --- fallback (very naive) ---
    try:
        text = Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        _log(f"ERROR parsing {path}: {e}")
        return []

    # extract sections that look like items (very permissive)
    blocks = re.split(r"\n{2,}|<hr[^>]*>", text, flags=re.I)
    for b in blocks:
        title = ""
        m = re.search(r"(?:^|\n)\s*Title\s*:\s*(.+)", b, flags=re.I)
        if m:
            title = m.group(1).strip()
        else:
            # fallback: first non-empty line
            lines = [ln.strip() for ln in re.split(r"<[^>]+>|\r|\n", b) if ln.strip()]
            title = lines[0][:140] if lines else ""

        if not title:
            continue

        item: Dict = {
            "title": title,
            "summary": "",
            "description": "",
            "status": "",
            "roadmap_id": "",
            "products": [],
            "platforms": [],
            "clouds": [],
            "ga": "",
            "_source": src,
        }

        def grab(label: str) -> str:
            m = re.search(rf"{label}\s*:\s*(.+)", b, flags=re.I)
            return (m.group(1).strip() if m else "")

        item["status"] = grab("Status")
        item["roadmap_id"] = grab("Roadmap ID") or grab("ID")
        item["summary"] = grab("Summary") or grab("Description")
        item["description"] = item["summary"]
        prod = grab("Product") or grab("Products")
        if prod:
            item["products"] = [p.strip() for p in re.split(r",|/|;", prod) if p.strip()]
        plat = grab("Platform") or grab("Platforms")
        if plat:
            item["platforms"] = [p.strip() for p in re.split(r",|/|;", plat) if p.strip()]
        clouds = grab("Cloud") or grab("Clouds")
        if clouds:
            item["clouds"] = [p.strip() for p in re.split(r",|/|;", clouds) if p.strip()]
        item["ga"] = grab("GA") or grab("Release") or ""

        items.append(item)

    _log(f"Parsed {len(items)} with month='{month}' from {path}")
    return items

def _parse_inputs(inputs: List[str], month: str) -> List[Dict]:
    all_items: List[Dict] = []
    for p in inputs:
        all_items.extend(_parse_one(p, month))
    return all_items

# ---- merge & fuzzy dedupe (MC-first) ----------------------------------------
def _norm(s: str) -> str:
    s = (s or "").lower()
    s = re.sub(r"[^a-z0-9]+", " ", s)
    return re.sub(r"\s+", " ", s).strip()

def _token_set(s: str) -> set:
    return set(_norm(s).split()) if s else set()

def _similar(a: str, b: str) -> float:
    """Jaccard similarity over tokens of titles."""
    ta, tb = _token_set(a), _token_set(b)
    if not ta or not tb:
        return 0.0
    inter = len(ta & tb)
    union = len(ta | tb)
    return inter / union if union else 0.0

def _has_val(v) -> bool:
    if v is None: return False
    if isinstance(v, str): return v.strip() != ""
    if isinstance(v, (list, tuple, set)): return len(v) > 0
    return True

def _pick(preferred, fallback):
    return preferred if _has_val(preferred) else fallback

def _norm_title(t: str) -> str:
    t = (t or "").lower().strip()
    t = re.sub(r"\s+", " ", t)
    t = re.sub(r"[^a-z0-9 ]+", "", t)
    return t

def _group_key(it: Dict) -> str:
    rid = (it.get("roadmap_id") or "").strip()
    return f"id:{rid}" if rid else f"title:{_norm_title(it.get('title') or it.get('headline') or '')[:80]}"

def _is_mc(it: Dict) -> bool:
    src = (it.get("source") or it.get("src") or "").lower()
    if src:
        return "message" in src or src == "mc"
    url = (it.get("url") or "").lower()
    return "messagecenter" in url or "admin.microsoft.com" in url

def _merge_pair_mc_first(mc: Dict, rm: Dict) -> Dict:
    out: Dict = {}

    # Identity
    out["roadmap_id"] = _pick(mc.get("roadmap_id"), rm.get("roadmap_id"))
    out["title"]      = _pick(mc.get("title"),      rm.get("title"))
    out["url"]        = _pick(mc.get("url"),        rm.get("url"))

    out["required_license"] = _pick(mc.get("required_license"), rm.get("required_license"))
    out["impact"]           = _pick(mc.get("impact"),           rm.get("impact"))
    out["how_to_implement"] = _pick(mc.get("how_to_implement"), rm.get("how_to_implement"))
    out["rollout_start"]    = _pick(mc.get("rollout_start"),    rm.get("rollout_start"))
 
    # Narrative → prefer MC
    out["summary"]     = _pick(mc.get("summary"),     rm.get("summary"))
    out["description"] = _pick(mc.get("description"), rm.get("description"))

    # Metadata → prefer RM (MC often blank)
    out["status"]    = _pick(mc.get("status"),    rm.get("status"))
    out["phases"]    = _pick(mc.get("phases"),    rm.get("phases"))
    out["products"]  = _pick(mc.get("products"),  rm.get("products"))
    out["platforms"] = _pick(mc.get("platforms"), rm.get("platforms"))
    out["clouds"]    = _pick(mc.get("clouds"),    rm.get("clouds"))
    # rollout/GA dates (whatever is present)
    out["rollout_start"] = _pick(mc.get("rollout_start"), rm.get("rollout_start"))
    out["ga_start"]      = _pick(mc.get("ga_start"),      rm.get("ga_start"))
    out["ga_end"]        = _pick(mc.get("ga_end"),        rm.get("ga_end"))

    # Audience flags if you parse them later
    out["aud_end_users"] = _pick(mc.get("aud_end_users"), rm.get("aud_end_users"))
    out["aud_admins"]    = _pick(mc.get("aud_admins"),    rm.get("aud_admins"))

    out["source"] = "mc+rm" if (mc and rm) else ("mc" if mc else "rm")
    return out

def _merge_mc_first(all_items: List[Dict]) -> List[Dict]:
    buckets: Dict[str, Dict[str, Optional[Dict]]] = {}
    for it in all_items:
        key = _group_key(it)
        b = buckets.setdefault(key, {"mc": None, "rm": None, "any": it})
        (b.__setitem__("mc", it) if _is_mc(it) else b.__setitem__("rm", it))
    merged: List[Dict] = []
    for b in buckets.values():
        mc, rm = b["mc"], b["rm"]
        merged.append(_merge_pair_mc_first(mc, rm) if (mc and rm) else (mc or rm or b["any"]))
    return merged

# ---- right rail rendering helpers -------------------------------------------
def _add_rail(slide, rail_width_in: float = 3.5, hex_color: str = "0F172A"):
    """Draw the right vertical rail across full height."""
    if not rail_width_in or rail_width_in <= 0:
        return
    # sanitize hex (pptx expects 'RRGGBB')
    hex_color = (hex_color or "0F172A").replace("#", "")
    full_w, full_h = int(10 * 914400), int(7.5 * 914400)
    rail_w = int(rail_width_in * 914400)
    left   = full_w - rail_w

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, 0, rail_w, full_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    shp.line.fill.background()







def _rail_color_for(item: Optional[Dict], assets: Dict) -> str:
    """Pick a rail color from product_palette with very forgiving inputs."""
    palette = assets.get("product_palette") or {}
    if not isinstance(palette, dict):
        palette = {}

    products = (item or {}).get("products") or (item or {}).get("product") or []
    # Coerce to list[str]
    if isinstance(products, str):
        prods: Iterable[str] = [products]
    elif isinstance(products, (list, tuple, set)):
        prods = [str(p) for p in products]
    else:
        prods = []

    key = (prods[0] if prods else "").strip().lower()
    color = palette.get(key) or palette.get("default") or "0F172A"
    return str(color).replace("#", "")


def _status_icon_for(status: str, assets: Dict) -> Optional[str]:
    s = (status or "").lower()
    rocket  = assets.get("icon_rocket")
    preview = assets.get("icon_preview")
    dev     = assets.get("icon_dev")

    if any(k in s for k in ("launched", "rolling out", "ga", "general availability", "available")):
        return rocket if rocket and Path(rocket).exists() else None
    if "preview" in s:
        return preview if preview and Path(preview).exists() else None
    if "develop" in s or "in dev" in s:
        return dev if dev and Path(dev).exists() else (preview if preview and Path(preview).exists() else None)
    return None




def _status_icon(slide, rail_left_in: float, rail_w_in: float, item: Dict, assets: Dict) -> None:
    path = _status_icon_for(item.get("status"), assets)
    if not path:
        return
    ICON_W, ICON_H = 0.90, 0.90
    left = rail_left_in + rail_w_in - ICON_W - 0.35
    top  = 0.10
    try:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 width=Inches(ICON_W), height=Inches(ICON_H))
    except Exception:
        pass


def _short_status_label(s: str) -> str:
    s = (s or "").strip().lower()
    if "launch" in s or "ga" in s or "available" in s:
        return "GA"
    if "rolling" in s:
        return "Rollout"
    if "develop" in s or "preview" in s:
        return "In Dev"
    return "Status"


def _status_chip(slide, rail_left_in: float, rail_w_in: float, status_text: str):
    """Compact status pill centered in the rail."""
    CHIP_W_IN, CHIP_H_IN = 2.2, 0.62
    CHIP_FONT_PT = 16
    left = rail_left_in + (rail_w_in - CHIP_W_IN) / 2.0
    top  = 0.50

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(CHIP_W_IN), Inches(CHIP_H_IN))
    shp.adjustments[0] = 0.25
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string("495057")
    shp.line.fill.background()

    tf = shp.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = _short_status_label(status_text)
    r.font.size = Pt(CHIP_FONT_PT); r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("FFFFFF")
    p.alignment = PP_ALIGN.CENTER

def _normalize_audience(item: dict) -> Tuple[Optional[bool], Optional[bool]]:
    """
    Returns (end_users, admins) as booleans or None when unknown.
    Accepts dict/list/string forms from parsers.
    """
    aud = item.get("audience")
    eu = ad = None

    def yesish(v: str) -> bool:
        v = (v or "").strip().lower()
        return v in ("y","yes","true","1","✓","✔","☑") or "yes" in v

    if isinstance(aud, dict):
        if "end_users" in aud: eu = yesish(str(aud.get("end_users")))
        if "end user" in aud:  eu = yesish(str(aud.get("end user")))
        if "admins" in aud:    ad = yesish(str(aud.get("admins")))
        if "admin" in aud:     ad = yesish(str(aud.get("admin")))
    elif isinstance(aud, list):
        low = [str(x).lower() for x in aud]
        eu = True  if any("end user" in x or x == "users" for x in low) else eu
        ad = True  if any("admin" in x for x in low) else ad
    elif isinstance(aud, str):
        s = aud.lower()
        if any(k in s for k in ("end user","end-user","users")): eu = True
        if "admin" in s: ad = True

    return (eu, ad)

def _audience_icons(slide, rail_left_in: float, rail_w_in: float, assets: Dict,
                    card_top_in: float, card_h_in: float,
                    item: Dict) -> None:
    left = rail_left_in + 0.25
    row_w = max(rail_w_in - 0.5, 1.2)
    top  = card_top_in + card_h_in + 0.35
    icon_w = icon_h = 0.5
    gap = 0.30

    icons = []
    if assets.get("icon_endusers") and Path(assets["icon_endusers"]).exists():
        icons.append(assets["icon_endusers"])
    if assets.get("icon_admins") and Path(assets["icon_admins"]).exists():
        icons.append(assets["icon_admins"])

    if not icons:
        return  # nothing to render

    total_w = len(icons) * icon_w + (len(icons)-1) * gap
    x = left + max((row_w - total_w)/2.0, 0.0)
    for p in icons:
        slide.shapes.add_picture(p, Inches(x), Inches(top), width=Inches(icon_w), height=Inches(icon_h))
        x += icon_w + gap






def _rail_audience_rows(slide, rail_left_in: float, rail_w_in: float, assets: Dict,
                        eu: Optional[bool], ad: Optional[bool],
                        card_top_in: float, card_h_in: float,
                        item: Dict) -> None:
    """Header + two rows placed BELOW the rail card, then optional Clouds + URL if space remains."""
    if rail_w_in <= 0:
        return

    left   = rail_left_in + 0.25
    row_w  = max(rail_w_in - 0.5, 1.4)
    row_h  = 0.55

    # Start just below the card
    start_top = card_top_in + card_h_in + 0.30

    # Header
    hdr = slide.shapes.add_textbox(Inches(left), Inches(start_top),
                                   Inches(row_w), Inches(0.5)).text_frame
    p = hdr.paragraphs[0]; r = p.add_run()
    r.text = "Target Audience:"
    r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("CBD5E1")

    rows_top = start_top + 0.60
    rows = [
        ("End Users:", assets.get("icon_endusers"), eu),
        ("Admins:",    assets.get("icon_admins"),   ad),
    ]

    for idx, (label, icon, val) in enumerate(rows):
        top = rows_top + idx * (row_h + 0.20)
        if icon and Path(icon).exists():
            slide.shapes.add_picture(icon, Inches(left), Inches(top),
                                     width=Inches(0.42), height=Inches(0.42))

        t = slide.shapes.add_textbox(Inches(left + 0.55), Inches(top),
                                     Inches(row_w - 0.55), Inches(row_h)).text_frame
        p = t.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{label} "; r1.font.bold = True
        r1.font.size = Pt(14); r1.font.color.rgb = RGBColor.from_string("FFFFFF")
        r2 = p.add_run()
        val_txt = "Yes" if val is True else "?" if val is None else "No"
        r2.text = val_txt; r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor.from_string("E5E7EB")

    # ---------- Optional extras UNDER Admin row (Clouds + Roadmap link) ----------
    extras_top = rows_top + 2 * (row_h + 0.20) + 0.25
    bottom_margin = 7.5 - 0.5  # slide_h - margin
    room_in = bottom_margin - extras_top

    if room_in >= 0.7:
        # Clouds
        clouds = item.get("clouds")
        clouds_text = ", ".join(clouds) if isinstance(clouds, (list, tuple, set)) else (clouds or "")
        if clouds_text:
            t = slide.shapes.add_textbox(Inches(left + 0.55), Inches(extras_top),
                                         Inches(row_w - 0.55), Inches(0.5)).text_frame
            p = t.paragraphs[0]
            r1 = p.add_run(); r1.text = "Clouds: "; r1.font.bold = True
            r1.font.size = Pt(13); r1.font.color.rgb = RGBColor.from_string("FFFFFF")
            r2 = p.add_run(); r2.text = clouds_text
            r2.font.size = Pt(13); r2.font.color.rgb = RGBColor.from_string("E5E7EB")
            extras_top += 0.55

    if room_in >= 1.2:
        # Roadmap link (as a small hyperlink)
        url = item.get("url") or item.get("link") or ""
        if url:
            tb = slide.shapes.add_textbox(Inches(left + 0.55), Inches(extras_top),
                                          Inches(row_w - 0.55), Inches(0.45))
            tf = tb.text_frame; p = tf.paragraphs[0]
            r = p.add_run(); r.text = "Open Roadmap ↗"
            r.font.size = Pt(12); r.font.color.rgb = RGBColor.from_string("93C5FD")
            try:
                r.hyperlink.address = url
            except Exception:
                pass


# ---- layout constants --------------------------------------------------------
# ---- layout constants --------------------------------------------------------
LINE_HEIGHT_IN, MIN_TITLE_H, MAX_TITLE_H = 0.55, 1.05, 1.40  # clamp title box tighter
TITLE_FONT_PT = 36                                           # was 44
BODY_FONT_PT  = 18                                           # was 20
TITLE_TOP     = 0.80


def _estimate_title_height(text: str, width_in: float) -> float:
    """Crude estimate based on characters-per-line; clamps to MIN/MAX."""
    if not text:
        return MIN_TITLE_H
    # assume ~ 36 chars per 6.0in width, scale linearly
    cpp = 36 * (width_in / 6.0)
    lines = max(1, int((len(text) / max(1, cpp)) + 0.5))
    h = lines * (LINE_HEIGHT_IN * 0.92)
    return max(MIN_TITLE_H, min(MAX_TITLE_H, h))

# ---- slide builder -----------------------------------------------------------
# before:
# def _build_item_slide(prs, item: Dict, month: str, assets: Dict, rail_width) -> None:
# after:
def _build_item_slide(prs, item: Dict, month: str, assets: Dict, rail_width,
                      page: int | None = None, total: int | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rail_w   = float(rail_width) if rail_width not in (None, "") else 0.0
    rail_left = 10.0 - rail_w
    w_body    = 10.0 - (rail_w + 1.2)

    # background
    S.add_full_slide_picture(slide, prs, assets.get("brand_bg") or assets.get("cover"))
    _add_rail(slide, rail_width_in=float(rail_width) if rail_width else 0.0,
                hex_color=_rail_color_for(item, assets))


    # status chip + icon
    status_text = (item.get("status") or "In Development").title()
    _status_chip(slide, rail_left, rail_w, status_text)
    _status_icon(slide, rail_left, rail_w, item, assets)

    # rail card (Month, Roadmap ID, Tenant, GA)
    card_top_in, card_h_in = _rail_card(slide, rail_left, rail_w, month, item)

    # audience rows (below the card)
    eu, ad = _normalize_audience(item)
   # target audience + optional clouds/url
    eu, ad = _normalize_audience(item)
    _rail_audience_rows(slide, rail_left, rail_w, assets, eu, ad, card_top_in, card_h_in, item)

    # left side content: centered title + summary paragraph
    title = item.get("title","") or item.get("headline","")

    body = (item.get("summary") or item.get("description") or "").strip()
    sections = []
    if item.get("required_license"):
        sections.append(f"Required license? {item['required_license']}")
    if item.get("impact"):
        sections.append(f"What is impact of this change? {item['impact']}")
    if item.get("how_to_implement"):
        sections.append(f"How to implement it {item['how_to_implement']}")

    if sections:
        body = (body + ("\n\n" if body else "")) + "\n\n".join(sections)



    title_h = _estimate_title_height(title, w_body)
    S.add_title_box(
        slide, title,
        left_in=0.6, top_in=TITLE_TOP, width_in=w_body, height_in=title_h,
        font_size_pt=TITLE_FONT_PT, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER
    )
    S.add_text_box(
        slide, body,
        left_in=0.6, top_in=2.0, width_in=8.8, height_in=3.6,
        font_size_pt=20, color="E5E7EB", auto_height=True
    )

    # footer
    footer_txt = month or ""
    if page is not None and total is not None:
        footer_txt = f"{footer_txt}  •  Page {page} of {total}".strip()


        S.add_text_box(
        slide, footer_txt,
        left_in=0.6, top_in=7.10, width_in=9.4, height_in=0.35,
        font_size_pt=12, color="94A3B8"
    )

# ---- public build() ----------------------------------------------------------
def build(
    inputs: List[str],
    output_path: str,
    month: str,
    assets: Dict,
    template: Optional[str] = None,
    rail_width: Optional[float] = None,
    conclusion_links: Optional[List[Dict]] = None,
    debug_dump: Optional[str] = None
) -> None:
    _log("Starting deck build...")

    # parse items
    all_items = _parse_inputs(inputs, month)
    if not all_items and month:
        _log("No items after month filter — retrying with no month filter...")
        all_items = _parse_inputs(inputs, "")

    if debug_dump:
        try:
            Path(debug_dump).write_text(json.dumps(all_items, indent=2), encoding="utf-8")
            _log(f"Debug dump wrote: {debug_dump}")
        except Exception as e:
            _log(f"Debug dump error: {e}")

    _log(f"Raw items before merge: {len(all_items)}")
    items = _merge_mc_first(all_items)
    _log(f"After MC-first merge + fuzzy: {len(items)} unique items")

    # presentation
    prs = Presentation(template) if (template and Path(template).exists()) else Presentation()

    # cover slide (if cover image present)
    if assets.get("cover") and Path(assets["cover"]).exists():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        S.add_full_slide_picture(slide, prs, assets["cover"])
        # optional logos
        if assets.get("logo") and Path(assets["logo"]).exists():
            slide.shapes.add_picture(assets["logo"], Inches(0.6), Inches(0.6), width=Inches(1.6))
        if assets.get("logo2") and Path(assets["logo2"]).exists():
            slide.shapes.add_picture(assets["logo2"], Inches(2.4), Inches(0.6), width=Inches(1.6))
        # title/date
        S.add_title_box(slide, assets.get("cover_title") or "M365 Technical Update Briefing",
                        left_in=0.6, top_in=2.2, width_in=8.8, height_in=1.6,
                        font_size_pt=50, bold=True, color="FFFFFF", align=PP_ALIGN.LEFT)
        S.add_text_box(slide, assets.get("cover_dates") or month or "",
                       left_in=0.6, top_in=3.9, width_in=8.8, height_in=0.8,
                       font_size_pt=24, color="E5E7EB")

    # agenda slide
    if assets.get("agenda") and Path(assets["agenda"]).exists():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        S.add_full_slide_picture(slide, prs, assets["agenda"])

    # item slides
    # assume 'items' (or whatever your final list variable is) contains the unique items to render
    total = len(items)
    for idx, it in enumerate(items, start=1):
        _build_item_slide(prs, it, month, assets, rail_width, page=idx, total=total)


    # conclusion/thank you
    if assets.get("conclusion") and Path(assets["conclusion"]).exists():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        S.add_full_slide_picture(slide, prs, assets["conclusion"])
    if assets.get("thankyou") and Path(assets["thankyou"]).exists():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        S.add_full_slide_picture(slide, prs, assets["thankyou"])

    # save
    out_abs = str(Path(output_path).resolve())
    prs.save(out_abs)
    _log(f"Deck saved: {out_abs}")

# ---- CLI ---------------------------------------------------------------------
if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("-i", "--inputs", nargs="+", required=True, help="HTML inputs")
    ap.add_argument("-o", "--output", required=True, help="Output PPTX")
    ap.add_argument("--month", default="", help="Month label for slides")
    ap.add_argument("--template", default="", help="Optional PPTX template path")

    # assets (underscore style for run_build)
    ap.add_argument("--cover", default="", help="Cover background image")
    ap.add_argument("--agenda", default="", help="Agenda background image")
    ap.add_argument("--separator", default="", help="Separator background image")
    ap.add_argument("--conclusion", default="", help="Conclusion background image")
    ap.add_argument("--thankyou", default="", help="Thank-you background image")
    ap.add_argument("--brand_bg", default="", help="Brand background for item slides")
    ap.add_argument("--cover_title", default="M365 Technical Update Briefing", help="Cover title")
    ap.add_argument("--cover_dates", default="", help="Cover date text")
    ap.add_argument("--logo", default="", help="Primary logo")
    ap.add_argument("--logo2", default="", help="Secondary logo")

    # icons
    ap.add_argument("--icon_rocket", default="", help="Icon for GA/Rollout")
    ap.add_argument("--icon_preview", default="", help="Icon for Preview/In Dev")
    ap.add_argument("--icon_dev", default="", help="Icon for In Development")

    ap.add_argument("--icon_endusers", default="", help="Icon for End Users audience")
    ap.add_argument("--icon_admins", default="", help="Icon for Admins audience")

    ap.add_argument("--rail_width", type=float, default=None, help="Right rail width in inches (e.g., 3.5)")
    ap.add_argument("--debug_dump", default="", help="Write parsed items JSON here")

    args = ap.parse_args()

    # resolve assets
    def g(p: str) -> str:
        return p if p and Path(p).exists() else ""

    assets = {
        "cover":      g(args.cover),
        "agenda":     g(args.agenda),
        "separator":  g(args.separator),
        "conclusion": g(args.conclusion),
        "thankyou":   g(args.thankyou),
        "brand_bg":   g(args.brand_bg),
        "cover_title": args.cover_title,
        "cover_dates": args.cover_dates or args.month,
        "logo": g(args.logo),
        "logo2": g(args.logo2),
        # icons
        "icon_rocket":   g(args.icon_rocket),
        "icon_preview":  g(args.icon_preview),
        "icon_endusers": g(args.icon_endusers),
        "icon_admins":   g(args.icon_admins),
        # default palette (can be overridden by generate_deck.py/style)
        "product_palette": {
            "teams": "4F46E5",
            "sharepoint": "16A34A",
            "onedrive": "0EA5E9",
            "exchange": "F97316",
            "outlook": "2563EB",
            "default": "0F172A",
        },
    }

    build(
        inputs=args.inputs,
        output_path=args.output,
        month=args.month,
        assets=assets,
        template=args.template,
        rail_width=args.rail_width,
        conclusion_links=None,
        debug_dump=(args.debug_dump or None),
    )
